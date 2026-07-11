#!/usr/bin/env python3
"""Unified media format converter — images, video/audio, and slide decks."""

from __future__ import annotations

import argparse
import json
import re
import shlex
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from arka.core.compute import ffmpeg_thread_args
from arka.media.compose_slides import (
    FORMAT_EXTENSIONS as SLIDE_EXT,
    SLIDE_FORMATS,
    _EXPORTERS,
    _build_json_export,
    _build_pdf,
    _metadata_payload,
    normalize_format as normalize_slide_format,
)
from arka.media.compose_video import Scene, _require_ffmpeg, _which

MediaType = str  # image | video | audio | slides

IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff", ".ico", ".svg", ".heic", ".heif"})
VIDEO_EXTS = frozenset({".mp4", ".webm", ".mov", ".avi", ".mkv", ".m4v", ".flv", ".wmv"})
AUDIO_EXTS = frozenset({".mp3", ".wav", ".aac", ".m4a", ".flac", ".ogg", ".opus"})
SLIDE_EXTS = frozenset({".pptx", ".ppt", ".pdf", ".html", ".htm", ".md", ".json"})

IMAGE_FORMATS = ("png", "jpg", "jpeg", "webp", "gif", "bmp", "tiff", "tif", "ico")
VIDEO_FORMATS = ("mp4", "webm", "mov", "avi", "mkv", "gif")
AUDIO_FORMATS = ("mp3", "wav", "aac", "m4a", "flac", "ogg", "opus")

IMAGE_ALIASES = {
    "png": "png",
    "jpg": "jpg",
    "jpeg": "jpg",
    "webp": "webp",
    "gif": "gif",
    "bmp": "bmp",
    "tiff": "tiff",
    "tif": "tiff",
    "ico": "ico",
    "svg": "svg",
    "heic": "heic",
    "heif": "heic",
}
VIDEO_ALIASES = {
    "mp4": "mp4",
    "webm": "webm",
    "mov": "mov",
    "avi": "avi",
    "mkv": "mkv",
    "gif": "gif",
    "m4v": "mp4",
}
AUDIO_ALIASES = {
    "mp3": "mp3",
    "wav": "wav",
    "aac": "aac",
    "m4a": "m4a",
    "flac": "flac",
    "ogg": "ogg",
    "opus": "opus",
}

PILLOW_SAVE = {
    "png": {"format": "PNG"},
    "jpg": {"format": "JPEG", "quality": 90},
    "jpeg": {"format": "JPEG", "quality": 90},
    "webp": {"format": "WEBP", "quality": 85},
    "gif": {"format": "GIF"},
    "bmp": {"format": "BMP"},
    "tiff": {"format": "TIFF"},
    "tif": {"format": "TIFF"},
    "ico": {"format": "ICO"},
}

_MAGIC = (
    (b"\x89PNG\r\n\x1a\n", ".png", "image"),
    (b"\xff\xd8\xff", ".jpg", "image"),
    (b"GIF87a", ".gif", "image"),
    (b"GIF89a", ".gif", "image"),
    (b"RIFF", ".webp", "image"),  # needs WEBP at offset 8
    (b"%PDF", ".pdf", "slides"),
    (b"\x00\x00\x00\x18ftyp", ".mp4", "video"),
    (b"\x00\x00\x00\x1cftyp", ".mp4", "video"),
    (b"\x00\x00\x00\x20ftyp", ".mp4", "video"),
    (b"ID3", ".mp3", "audio"),
    (b"\xff\xfb", ".mp3", "audio"),
    (b"PK\x03\x04", ".pptx", "slides"),
)


def _normalize_image_format(name: str) -> str:
    return IMAGE_ALIASES.get((name or "").strip().lower().lstrip("."), "")


def _normalize_video_format(name: str) -> str:
    return VIDEO_ALIASES.get((name or "").strip().lower().lstrip("."), "")


def _normalize_audio_format(name: str) -> str:
    return AUDIO_ALIASES.get((name or "").strip().lower().lstrip("."), "")


def detect_media_type(path: Path) -> MediaType:
    """Detect image | video | audio | slides from extension and magic bytes."""
    path = path.expanduser()
    ext = path.suffix.lower()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in VIDEO_EXTS:
        return "video"
    if ext in AUDIO_EXTS:
        return "audio"
    if ext in SLIDE_EXTS:
        return "slides"
    if path.is_file():
        head = path.read_bytes()[:32]
        if len(head) >= 12 and head[:4] == b"RIFF" and head[8:12] == b"WEBP":
            return "image"
        for magic, guess_ext, media_type in _MAGIC:
            if head.startswith(magic):
                if guess_ext == ".pdf":
                    return "slides"
                return media_type
        if head[:4] == b"PK\x03\x04":
            return "slides"
    raise SystemExit(f"Cannot detect media type for {path.name!r} — unsupported extension {ext!r}")


def _formats_for_type(media_type: MediaType) -> tuple[str, ...]:
    if media_type == "image":
        return IMAGE_FORMATS
    if media_type == "video":
        return VIDEO_FORMATS + AUDIO_FORMATS
    if media_type == "audio":
        return AUDIO_FORMATS
    return SLIDE_FORMATS


def _normalize_target_format(name: str, media_type: MediaType) -> str:
    raw = (name or "").strip().lower().lstrip(".")
    if raw == "all":
        return "all"
    if media_type == "image":
        return _normalize_image_format(raw)
    if media_type == "audio":
        return _normalize_audio_format(raw)
    if media_type == "video":
        fmt = _normalize_video_format(raw) or _normalize_audio_format(raw)
        return fmt
    return normalize_slide_format(raw)


def parse_target_formats(value: str | None, media_type: MediaType) -> list[str]:
    raw = (value or "").strip().lower()
    if not raw:
        raise SystemExit("Target format required — use --to FORMAT or -o output.ext")
    if raw == "all":
        return list(_formats_for_type(media_type))
    formats: list[str] = []
    for token in re.split(r"[,+\s]+", raw):
        token = token.strip().lstrip(".")
        if not token:
            continue
        fmt = _normalize_target_format(token, media_type)
        if not fmt:
            known = ", ".join(_formats_for_type(media_type) + ("all",))
            raise SystemExit(f"Unknown {media_type} format {token!r}. Choose: {known}")
        if fmt not in formats:
            formats.append(fmt)
    return formats


def _default_output(input_path: Path, fmt: str, media_type: MediaType) -> Path:
    stem = input_path.with_suffix("")
    if media_type == "slides":
        ext = SLIDE_EXT.get(fmt, f".{fmt}")
        return stem.with_suffix(ext)
    return stem.with_suffix(f".{fmt}")


def _output_paths(input_path: Path, formats: list[str], media_type: MediaType, explicit: Path | None) -> dict[str, Path]:
    if explicit and len(formats) == 1:
        return {formats[0]: explicit}
    if explicit and len(formats) > 1:
        stem = explicit.with_suffix("")
        return {
            fmt: (stem.with_suffix(SLIDE_EXT[fmt]) if media_type == "slides" else stem.with_suffix(f".{fmt}"))
            for fmt in formats
        }
    return {fmt: _default_output(input_path, fmt, media_type) for fmt in formats}


def _register_heif() -> bool:
    try:
        from pillow_heif import register_heif_opener

        register_heif_opener()
        return True
    except ImportError:
        return False


def _require_pillow_image():
    try:
        from PIL import Image
    except ImportError as exc:
        raise SystemExit(
            "Pillow is required for image conversion.\nInstall: pip install Pillow  or  pip install 'arka-agent[drawings]'"
        ) from exc
    return Image


def _open_image(path: Path):
    Image = _require_pillow_image()
    if path.suffix.lower() in {".heic", ".heif"} and not _register_heif():
        raise SystemExit(
            "HEIC input requires pillow-heif.\nInstall: pip install pillow-heif  or  pip install 'arka-agent[drawings]'"
        )
    if path.suffix.lower() == ".svg":
        try:
            import cairosvg
        except ImportError as exc:
            raise SystemExit(
                "SVG input requires cairosvg.\nInstall: pip install cairosvg  (also needs Cairo on the system)"
            ) from exc
        png_bytes = cairosvg.svg2png(url=str(path))
        from io import BytesIO

        return Image.open(BytesIO(png_bytes))
    return Image.open(path)


def convert_image(
    input_path: Path,
    output_path: Path,
    *,
    target_format: str,
    quality: int | None = None,
    width: int | None = None,
    height: int | None = None,
) -> Path:
    Image = _require_pillow_image()
    img = _open_image(input_path)
    fmt = _normalize_image_format(target_format) or target_format
    if fmt in {"jpg", "jpeg", "bmp", "ico"} and img.mode in {"RGBA", "LA", "P"}:
        bg = Image.new("RGB", img.size, (255, 255, 255))
        if img.mode == "P":
            img = img.convert("RGBA")
        bg.paste(img, mask=img.split()[-1] if img.mode in {"RGBA", "LA"} else None)
        img = bg
    elif img.mode not in {"RGB", "L", "RGBA", "P"}:
        img = img.convert("RGB")

    if width or height:
        w = width or img.width
        h = height or img.height
        if width and not height:
            h = int(img.height * (width / img.width))
        elif height and not width:
            w = int(img.width * (height / img.height))
        img = img.resize((w, h), Image.Resampling.LANCZOS)

    save_kwargs = dict(PILLOW_SAVE.get(fmt, {"format": fmt.upper()}))
    if quality is not None and save_kwargs.get("format") in {"JPEG", "WEBP"}:
        save_kwargs["quality"] = quality
    output_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(output_path), **save_kwargs)
    return output_path


def _ffmpeg_run(cmd: list[str]) -> None:
    proc = subprocess.run(cmd, capture_output=True, text=True, check=False)
    if proc.returncode != 0:
        raise SystemExit(f"ffmpeg failed: {(proc.stderr or proc.stdout or proc.returncode).strip()}")


def convert_video(
    input_path: Path,
    output_path: Path,
    *,
    target_format: str,
    trim_start: float | None = None,
    trim_duration: float | None = None,
) -> Path:
    ffmpeg = _require_ffmpeg()
    fmt = _normalize_video_format(target_format) or _normalize_audio_format(target_format) or target_format
    output_path.parent.mkdir(parents=True, exist_ok=True)

    ss_args: list[str] = []
    if trim_start is not None and trim_start > 0:
        ss_args = ["-ss", f"{trim_start:.3f}"]
    t_args: list[str] = []
    if trim_duration is not None and trim_duration > 0:
        t_args = ["-t", f"{trim_duration:.3f}"]

    if fmt == "gif":
        palette = output_path.with_suffix(".palette.png")
        try:
            _ffmpeg_run(
                [
                    ffmpeg,
                    "-nostdin",
                    "-hide_banner",
                    "-loglevel",
                    "error",
                    *ffmpeg_thread_args(),
                    "-y",
                    *ss_args,
                    "-i",
                    str(input_path),
                    *t_args,
                    "-vf",
                    "fps=12,scale=640:-1:flags=lanczos,palettegen",
                    str(palette),
                ]
            )
            _ffmpeg_run(
                [
                    ffmpeg,
                    "-nostdin",
                    "-hide_banner",
                    "-loglevel",
                    "error",
                    *ffmpeg_thread_args(),
                    "-y",
                    *ss_args,
                    "-i",
                    str(input_path),
                    "-i",
                    str(palette),
                    *t_args,
                    "-lavfi",
                    "fps=12,scale=640:-1:flags=lanczos[x];[x][1:v]paletteuse",
                    str(output_path),
                ]
            )
        finally:
            palette.unlink(missing_ok=True)
        return output_path

    if fmt in AUDIO_FORMATS:
        codec_map = {"mp3": "libmp3lame", "aac": "aac", "m4a": "aac", "wav": "pcm_s16le", "flac": "flac", "ogg": "libvorbis", "opus": "libopus"}
        cmd = [
            ffmpeg,
            "-nostdin",
            "-hide_banner",
            "-loglevel",
            "error",
            *ffmpeg_thread_args(),
            "-y",
            *ss_args,
            "-i",
            str(input_path),
            *t_args,
            "-vn",
            "-c:a",
            codec_map.get(fmt, "copy"),
            str(output_path),
        ]
        _ffmpeg_run(cmd)
        return output_path

    vcodec = {"mp4": "libx264", "webm": "libvpx-vp9", "mov": "libx264", "avi": "libx264", "mkv": "libx264"}.get(fmt, "libx264")
    cmd = [
        ffmpeg,
        "-nostdin",
        "-hide_banner",
        "-loglevel",
        "error",
        *ffmpeg_thread_args(),
        "-y",
        *ss_args,
        "-i",
        str(input_path),
        *t_args,
        "-c:v",
        vcodec,
        "-c:a",
        "aac" if fmt in {"mp4", "mov", "mkv"} else "copy",
        "-pix_fmt",
        "yuv420p",
        "-movflags",
        "+faststart",
        str(output_path),
    ]
    _ffmpeg_run(cmd)
    return output_path


def _load_arka_slide_json(path: Path) -> tuple[list[Scene], str, list[dict]]:
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise SystemExit("Slide JSON must be an object with scenes.")
    rows = data.get("scenes") or []
    topic = str(data.get("topic") or "slides")
    credits = data.get("photo_credits") or data.get("unsplash_credits") or []
    scenes: list[Scene] = []
    for row in rows:
        if not isinstance(row, dict):
            continue
        title = str(row.get("title") or "").strip()
        if not title:
            continue
        scenes.append(
            Scene(
                title=title,
                narration=str(row.get("narration") or ""),
                body=str(row.get("body") or ""),
                captions=[str(c) for c in (row.get("captions") or []) if str(c).strip()],
                image_query=str(row.get("image_query") or ""),
                image_keywords=[str(k) for k in (row.get("image_keywords") or []) if str(k).strip()],
                chart=row.get("chart") if isinstance(row.get("chart"), dict) else None,
                chart_path=str(row.get("chart_path") or ""),
                slide_image=str(row.get("slide_image") or ""),
            )
        )
    if not scenes:
        raise SystemExit("No scenes found in slide JSON.")
    return scenes, topic, credits if isinstance(credits, list) else []


def _extract_pptx_images(input_path: Path, work: Path) -> list[Path]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit(
            "python-pptx is required for PPTX conversion.\nInstall: pip install python-pptx  or  pip install 'arka-agent[video]'"
        ) from exc

    prs = Presentation(str(input_path))
    images: list[Path] = []
    for idx, slide in enumerate(prs.slides):
        out = work / f"slide-{idx:02d}.png"
        saved = False
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                img = shape.image
                out.write_bytes(img.blob)
                images.append(out)
                saved = True
                break
        if not saved:
            from arka.media.compose_video import load_config
            from arka.media.chart_slide import render_title_slide

            render_title_slide(Scene(title=f"Slide {idx + 1}"), out, load_config())
            images.append(out)
    if not images:
        raise SystemExit("No slides found in PPTX.")
    return images


def _pdf_pages_to_images(input_path: Path, work: Path) -> list[Path]:
    try:
        import fitz
    except ImportError as exc:
        raise SystemExit(
            "PDF slide conversion requires pymupdf.\nInstall: pip install pymupdf  or  pip install 'arka-agent[drawings]'"
        ) from exc
    doc = fitz.open(str(input_path))
    images: list[Path] = []
    for idx in range(len(doc)):
        page = doc.load_page(idx)
        pix = page.get_pixmap(dpi=150)
        out = work / f"slide-{idx:02d}.png"
        pix.save(str(out))
        images.append(out)
    doc.close()
    if not images:
        raise SystemExit("No pages found in PDF.")
    return images


def _slide_images_from_input(input_path: Path, work: Path) -> list[Path]:
    ext = input_path.suffix.lower()
    if ext in {".pptx", ".ppt"}:
        return _extract_pptx_images(input_path, work)
    if ext == ".pdf":
        return _pdf_pages_to_images(input_path, work)
    raise SystemExit(
        f"Slide conversion from {ext} is limited — supported inputs: pptx, pdf, arka JSON (.json with scenes)."
    )


def convert_slides(
    input_path: Path,
    output_paths: dict[str, Path],
    *,
    quality: int | None = None,
) -> list[Path]:
    _ = quality
    ext = input_path.suffix.lower()
    saved: list[Path] = []
    work = Path(tempfile.mkdtemp(prefix="arka-convert-slides-"))
    try:
        if ext == ".json":
            try:
                scenes, topic, credits = _load_arka_slide_json(input_path)
            except json.JSONDecodeError as exc:
                raise SystemExit(f"Invalid slide JSON: {exc}") from exc
            from arka.media.compose_slides import _render_scene_png, load_config
            from arka.media.compose_video import load_config as video_cfg

            cfg = load_config()
            slide_images: list[tuple[Path, Scene]] = []
            used: set[str] = set()
            credit_rows: list[dict] = list(credits)
            for i, scene in enumerate(scenes):
                if scene.slide_image:
                    png = Path(scene.slide_image).expanduser()
                    if png.is_file():
                        slide_images.append((png, scene))
                        continue
                png = _render_scene_png(
                    scene,
                    topic=topic,
                    work=work,
                    index=i,
                    cfg=cfg,
                    used_photo_ids=used,
                    credits=credit_rows,
                )
                slide_images.append((png, scene))
            outputs_map: dict[str, str] = {}
            for fmt, path in output_paths.items():
                if fmt == "json":
                    continue
                exporter = _EXPORTERS[fmt]
                if fmt in {"pptx", "pdf", "html"}:
                    saved_path = exporter(slide_images, path, topic=topic, cfg=video_cfg())
                else:
                    saved_path = exporter(slide_images, path, topic=topic, scenes=scenes, cfg=video_cfg())
                saved.append(saved_path)
                outputs_map[fmt] = str(saved_path)
            if "json" in output_paths:
                saved.append(
                    _build_json_export(
                        slide_images,
                        output_paths["json"],
                        topic=topic,
                        scenes=scenes,
                        credits=credit_rows,
                        outputs=outputs_map,
                        cfg=video_cfg(),
                    )
                )
            return saved

        images = _slide_images_from_input(input_path, work)
        slide_pairs = [(img, Scene(title=f"Slide {i + 1}")) for i, img in enumerate(images)]
        topic = input_path.stem

        for fmt, path in output_paths.items():
            if fmt == "json":
                payload = _metadata_payload(
                    topic=topic,
                    scenes=[s for _, s in slide_pairs],
                    credits=[],
                    outputs={},
                )
                path.parent.mkdir(parents=True, exist_ok=True)
                path.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")
                saved.append(path)
                continue
            if fmt == "pdf":
                saved.append(_build_pdf(slide_pairs, path))
                continue
            if fmt == "html":
                saved.append(_EXPORTERS["html"](slide_pairs, path, topic=topic))
                continue
            if fmt == "md":
                saved.append(
                    _EXPORTERS["md"](
                        slide_pairs,
                        path,
                        topic=topic,
                        scenes=[s for _, s in slide_pairs],
                    )
                )
                continue
            if fmt == "pptx":
                if ext in {".pptx", ".ppt"} and len(output_paths) == 1:
                    shutil.copy2(input_path, path)
                    saved.append(path)
                    continue
                from arka.media.compose_video import load_config as video_cfg

                saved.append(_EXPORTERS["pptx"](slide_pairs, path, topic=topic, cfg=video_cfg()))
                continue
            raise SystemExit(f"Cannot convert {input_path.suffix} → {fmt}")

        return saved
    finally:
        shutil.rmtree(work, ignore_errors=True)


def convert_media(
    input_path: Path,
    *,
    target_formats: list[str],
    output: Path | None = None,
    quality: int | None = None,
    width: int | None = None,
    height: int | None = None,
    trim_start: float | None = None,
    trim_duration: float | None = None,
) -> list[Path]:
    input_path = input_path.expanduser().resolve()
    if not input_path.is_file():
        raise SystemExit(f"Input not found: {input_path}")

    media_type = detect_media_type(input_path)
    paths = _output_paths(input_path, target_formats, media_type, output)
    saved: list[Path] = []

    if media_type == "image":
        for fmt, out in paths.items():
            saved.append(
                convert_image(
                    input_path,
                    out,
                    target_format=fmt,
                    quality=quality,
                    width=width,
                    height=height,
                )
            )
        return saved

    if media_type in {"video", "audio"}:
        for fmt, out in paths.items():
            saved.append(
                convert_video(
                    input_path,
                    out,
                    target_format=fmt,
                    trim_start=trim_start,
                    trim_duration=trim_duration,
                )
            )
        return saved

    return convert_slides(input_path, paths, quality=quality)


def _looks_like_currency(text: str) -> bool:
    try:
        from arka.integrations.currency import parse_convert

        return parse_convert(text) is not None
    except ImportError:
        return bool(re.search(r"(?i)\b(?:usd|eur|inr|gbp|dollars?|euros?|rupees?|pounds?)\b", text))


def nl_to_argv(text: str) -> list[str]:
    t = text.strip()
    if not t or _looks_like_currency(t):
        return []

    media_ext = r"(?:png|jpe?g|webp|gif|bmp|tiff?|ico|svg|heic|heif|mp4|webm|mov|avi|mkv|m4v|mp3|wav|aac|m4a|flac|ogg|opus|pptx?|pdf|html?|md|json)"
    fmt_token = r"[a-z0-9]{2,8}"

    patterns = [
        rf"(?i)^(?:please\s+)?(?:convert|export|transform)\s+(?P<input>\S+\.{media_ext})\s+(?:to|into|as)\s+(?P<fmt>{fmt_token}|all)\b",
        rf"(?i)^(?:please\s+)?(?:convert|export)\s+(?P<input>\S+\.{media_ext})\s+(?:to|into)\s+(?P<fmt>{fmt_token}|all)\s+format\b",
        rf"(?i)^(?:please\s+)?(?:convert|export)\s+(?:the\s+)?(?:video|image|photo|picture|slide|slides|deck|presentation)\s+(?P<input>\S+\.{media_ext})\s+(?:to|into|as)\s+(?P<fmt>{fmt_token}|all)\b",
        rf"(?i)^(?:please\s+)?(?:convert|export)\s+(?:my\s+)?(?:video|image|photo|picture|slide|slides|deck|presentation)\s+(?:to|into|as)\s+(?P<fmt>{fmt_token}|all)\b(?P<input>\s+\S+\.{media_ext})?",
    ]
    for pat in patterns:
        m = re.search(pat, t)
        if not m:
            continue
        input_file = (m.group("input") or "").strip()
        fmt = (m.group("fmt") or "").strip().lower()
        if not input_file:
            tail = re.search(rf"(\S+\.{media_ext})\s*$", t, re.I)
            if not tail:
                continue
            input_file = tail.group(1)
        if not fmt:
            continue
        argv = [input_file, "--to", fmt]
        return argv

    # "convert slides to pdf" with explicit path elsewhere
    m = re.search(
        rf"(?i)(?:convert|export)\s+(?P<input>\S+\.{media_ext})\s+(?:to|into)\s+(?P<fmt>pptx|pdf|html|markdown|md|json|all)\b",
        t,
    )
    if m:
        fmt = m.group("fmt").lower()
        if fmt == "markdown":
            fmt = "md"
        return [m.group("input"), "--to", fmt]
    return []


def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="Convert images, video/audio, and slide decks between formats")
    sub = p.add_subparsers(dest="cmd")

    p_convert = sub.add_parser("convert", help="Convert input file to target format(s)")
    p_convert.add_argument("input", help="Input file path")
    p_convert.add_argument("-t", "--to", "--format", dest="target", help="Target format (or all)")
    p_convert.add_argument("-o", "--output", help="Output path (single format only unless stem-based multi-export)")
    p_convert.add_argument("-q", "--quality", type=int, help="JPEG/WebP quality (1-100)")
    p_convert.add_argument("--width", type=int, help="Resize width (images)")
    p_convert.add_argument("--height", type=int, help="Resize height (images)")
    p_convert.add_argument("--trim-start", type=float, help="Trim start offset in seconds (video/audio)")
    p_convert.add_argument("--trim-duration", type=float, help="Trim duration in seconds (video/audio)")
    p_convert.set_defaults(func=cmd_convert)

    p_parse = sub.add_parser("parse", help="Parse natural language → convert args")
    p_parse.add_argument("text", nargs="+")
    p_parse.set_defaults(func=cmd_parse)

    p_check = sub.add_parser("check", help="Verify conversion tools")
    p_check.set_defaults(func=cmd_check)

    return p


def cmd_check(_args: argparse.Namespace) -> int:
    ok = True
    try:
        _require_pillow_image()
        print("✓ Pillow (images, slide PDF)")
    except SystemExit:
        print("✗ Pillow — pip install Pillow", file=sys.stderr)
        ok = False
    if _which("ffmpeg"):
        print("✓ ffmpeg (video/audio)")
    else:
        print("✗ ffmpeg — brew install ffmpeg", file=sys.stderr)
        ok = False
    try:
        import pptx  # noqa: F401

        print("✓ python-pptx (pptx slides)")
    except ImportError:
        print("  python-pptx optional — pip install python-pptx (pptx read/write)")
    try:
        import fitz  # noqa: F401

        print("✓ pymupdf (pdf slides)")
    except ImportError:
        print("  pymupdf optional — pip install pymupdf (pdf → images)")
    if _register_heif():
        print("✓ pillow-heif (HEIC)")
    else:
        print("  pillow-heif optional — pip install pillow-heif (HEIC)")
    if _which("soffice") or _which("libreoffice"):
        print("✓ LibreOffice (pptx fallback)")
    else:
        print("  LibreOffice optional — headless pptx conversion fallback")
    print("✓ html/md/json slide export (stdlib)")
    return 0 if ok else 1


def cmd_parse(args: argparse.Namespace) -> int:
    argv = nl_to_argv(" ".join(args.text))
    if not argv:
        return 1
    print(" ".join(shlex.quote(a) for a in argv))
    return 0


def cmd_convert(args: argparse.Namespace) -> int:
    input_path = Path(args.input).expanduser()
    media_type = detect_media_type(input_path)
    target = args.target
    if not target and args.output:
        target = args.output.suffix.lstrip(".")
    if not target:
        print("Provide --to FORMAT or -o output.ext", file=sys.stderr)
        return 1
    try:
        formats = parse_target_formats(target, media_type)
    except SystemExit as exc:
        print(exc, file=sys.stderr)
        return 1
    output = Path(args.output).expanduser() if args.output else None
    print(f"Converting {input_path.name} ({media_type}) → {', '.join(formats)}", file=sys.stderr)
    saved = convert_media(
        input_path,
        target_formats=formats,
        output=output,
        quality=args.quality,
        width=args.width,
        height=args.height,
        trim_start=args.trim_start,
        trim_duration=args.trim_duration,
    )
    for path in saved:
        print(f"Saved: {path}")
    return 0


def main(argv: list[str] | None = None) -> int:
    from arka.env import load_env

    load_env()
    argv = list(argv if argv is not None else sys.argv[1:])
    if not argv:
        build_parser().print_help()
        return 0
    if argv[0] not in {"convert", "parse", "check", "-h", "--help"}:
        if argv[0] == "check":
            argv = ["check"]
        elif Path(argv[0]).suffix:
            rest = argv[1:]
            nl = nl_to_argv(" ".join(argv))
            if nl:
                argv = ["convert", *nl]
            else:
                argv = ["convert", argv[0], *rest]
        else:
            nl = nl_to_argv(" ".join(argv))
            if nl:
                argv = ["convert", *nl]
            else:
                build_parser().print_help()
                return 1
    parser = build_parser()
    args = parser.parse_args(argv)
    if not getattr(args, "func", None):
        parser.print_help()
        return 0
    return int(args.func(args))


if __name__ == "__main__":
    raise SystemExit(main())
