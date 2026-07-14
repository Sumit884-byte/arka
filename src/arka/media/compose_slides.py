#!/usr/bin/env python3
"""Compose presentation slide decks — stock photos, charts, LLM scripts (like compose_video)."""

from __future__ import annotations

import argparse
import base64
import html as html_module
import io
import json
import os
import re
import shlex
import subprocess
import sys
import tempfile
import zipfile
from xml.etree import ElementTree as ET
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

from arka.media.compose_video import (
    Scene,
    VideoConfig,
    _attach_photo_queries,
    _enrich_scenes,
    _fetch_scene_photo,
    _llm_available,
    _parse_scenes_json,
    _render_photo_slide,
    _scene_search_query,
    _which,
    load_config,
    normalize_topic,
    topic_label,
)
from arka.media.stock_photos import any_source_available, setup_hint as stock_setup_hint

SLIDE_STYLES = ("executive", "academic", "pitch")
SLIDE_FORMATS = ("pptx", "pdf", "html", "md", "json")
FORMAT_EXTENSIONS = {
    "pptx": ".pptx",
    "pdf": ".pdf",
    "html": ".html",
    "md": ".md",
    "json": ".json",
}
FORMAT_ALIASES = {
    "pptx": "pptx",
    "ppt": "pptx",
    "powerpoint": "pptx",
    "pdf": "pdf",
    "html": "html",
    "htm": "html",
    "md": "md",
    "markdown": "md",
    "marp": "md",
    "json": "json",
}


def _env(name: str, default: str = "") -> str:
    return os.environ.get(name, default).strip()


def _env_int(name: str, default: int) -> int:
    raw = _env(name)
    if not raw:
        return default
    try:
        return int(raw)
    except ValueError:
        return default


def normalize_slide_style(name: str | None) -> str:
    raw = (name or _env("SLIDES_DEFAULT_STYLE", "executive")).strip().lower()
    return raw if raw in SLIDE_STYLES else "executive"


def _style_guidance(style: str) -> str:
    guides = {
        "executive": (
            "Executive deck: confident, concise, decision-oriented. "
            "Action titles state the insight (not section labels like 'Background'). "
            "Arc: hook → context → key insight → implications → recommendation → next steps. "
            "Tone: boardroom-ready, one idea per slide."
        ),
        "academic": (
            "Academic presentation: formal, precise, evidence-led. "
            "Titles name the claim or section clearly. "
            "Arc: intro → context/literature → core concepts → evidence/examples → implications → conclusion. "
            "Tone: seminar or conference talk; slightly more detail allowed."
        ),
        "pitch": (
            "Investor pitch deck: bold, urgent, outcome-focused. "
            "Action titles sell momentum (problem size, solution edge, traction). "
            "Arc: hook → problem → solution → market → traction → business model → team → ask/CTA. "
            "Tone: startup pitch; every slide earns the next."
        ),
    }
    return guides.get(normalize_slide_style(style), guides["executive"])


def extract_slide_style(text: str) -> str:
    t = text.strip().lower()
    if not t:
        return normalize_slide_style(None)
    for style in SLIDE_STYLES:
        if re.search(rf"\b{style}\b(?:\s+(?:style|deck|slides|presentation|tone))?", t):
            return style
        if re.search(rf"\b(?:style|tone)\s+{style}\b", t):
            return style
    return normalize_slide_style(None)


def _strip_style_from_text(text: str) -> str:
    t = text.strip()
    for style in SLIDE_STYLES:
        t = re.sub(rf"(?i)\b{style}\s+(?:style\s+)?", "", t)
        t = re.sub(rf"(?i)\b(?:style|tone)\s+{style}\b", "", t)
    return re.sub(r"\s{2,}", " ", t).strip(" ,;")


def _slides_scene_bounds() -> tuple[int, int]:
    min_s = max(3, _env_int("SLIDES_MIN_SCENES", 6))
    max_s = max(min_s, _env_int("SLIDES_MAX_SCENES", 12))
    return min_s, max_s


def _slides_script_mode(args: argparse.Namespace) -> str:
    if getattr(args, "llm", False):
        return "llm"
    mode = _env("SLIDES_COMPOSE_SCRIPT", "auto").lower()
    if mode in {"llm", "template"}:
        return mode
    return "llm" if _llm_available() else "template"


def _enrich_slide_scenes(scenes: list[Scene]) -> list[Scene]:
    from arka.media.compose_video import _caption_from_narration, _normalize_scene_text

    out: list[Scene] = []
    for scene in scenes:
        scene.narration = _normalize_scene_text(scene.narration)
        scene.body = _normalize_scene_text(scene.body)
        if scene.captions:
            scene.captions = [str(caption).strip() for caption in scene.captions if str(caption).strip()]
        caps = scene.captions[:4]
        if caps:
            scene.body = "\n".join(f"• {caption}" for caption in caps[:3])
        elif scene.body.strip():
            pass
        elif scene.narration.strip():
            scene.body = _caption_from_narration(scene.narration)
        out.append(scene)
    return out


def _slides_script_needs_shortening(scenes: list[Scene]) -> bool:
    _, max_scenes = _slides_scene_bounds()
    if len(scenes) > max_scenes:
        return True
    for scene in scenes:
        caps = scene.captions or []
        if len(caps) > 4:
            return True
        if any(len(str(caption).split()) > 12 for caption in caps):
            return True
        if len(scene.title.split()) > 14:
            return True
    return False


def _llm_slides_script(
    topic: str,
    *,
    scenes: int | None = None,
    style: str = "executive",
) -> list[Scene]:
    try:
        from arka.llm.fallback import llm_complete
    except ImportError as exc:
        raise SystemExit("LLM script generation requires arka chat deps (pip install 'arka-agent[chat]')") from exc

    min_scenes, max_scenes = _slides_scene_bounds()
    style = normalize_slide_style(style)
    system = (
        "You write high-quality presentation slide decks for live audiences. "
        "Return ONLY a JSON array (no markdown). Each item: "
        '{"title":"action title — insight, not a section label", '
        '"narration":"speaker notes (2-4 sentences, optional detail for presenter)", '
        '"body":"optional one-line subtitle on slide (max 10 words) OR empty string", '
        '"captions":["bullet 1","bullet 2"] (0-3 bullets, max 10 words each; omit on title-only slides), '
        '"image_keywords":["conference room","team whiteboard"] (3-5 short visual phrases, 2-3 words each), '
        '"image_query":"optional 2-4 word stock photo fallback", '
        '"chart":{"type":"bar|barh|pie|line|grouped_bar", "title":"...", '
        '"data":"Label:10,Other:20", "ylabel":"...", "source":"..."}} '
        "Rules: one idea per slide; no walls of text; titles are conclusions or claims; "
        "use charts only when numbers strengthen the story."
    )
    if scenes is not None:
        scene_hint = f"Slides: exactly {scenes}"
    else:
        scene_hint = (
            f"Choose an appropriate slide count ({min_scenes}-{max_scenes}) "
            "for a focused deck — not a document."
        )
    user = (
        f"Topic: {topic_label(topic)}\n"
        f"Style: {style}\n"
        f"{_style_guidance(style)}\n"
        f"{scene_hint}\n"
        "Keep on-screen text minimal; put depth in narration (speaker notes)."
    )
    text = llm_complete(system, user, temperature=0.35, task="compose_slides")
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    try:
        parsed = _parse_scenes_json(text)
    except (json.JSONDecodeError, ValueError) as exc:
        raise SystemExit(f"LLM returned invalid slide JSON: {exc}") from exc
    if len(parsed) < 2:
        raise SystemExit("LLM did not return a usable slide script.")
    if scenes is None and len(parsed) > max_scenes:
        parsed = parsed[:max_scenes]
    return _enrich_slide_scenes(parsed)


def _llm_slides_summarize_script(
    topic: str,
    scenes: list[Scene],
    *,
    style: str = "executive",
) -> list[Scene]:
    try:
        from arka.llm.fallback import llm_complete
    except ImportError as exc:
        raise SystemExit("LLM script generation requires arka chat deps (pip install 'arka-agent[chat]')") from exc

    payload = [
        {
            "title": scene.title,
            "narration": scene.narration,
            "body": scene.body,
            "captions": scene.captions,
            "image_query": scene.image_query,
            "image_keywords": scene.image_keywords,
            "chart": scene.chart,
        }
        for scene in scenes
    ]
    system = (
        "Tighten a presentation slide deck for on-screen readability. "
        "Return ONLY a JSON array with: title, narration, body, captions, image_query, chart. "
        "Keep action titles; shorten narration by 20-40%; max 3 captions per slide (max 10 words each); "
        "drop filler slides if the deck is too long."
    )
    user = (
        f"Topic: {topic_label(topic)}\n"
        f"Style: {normalize_slide_style(style)}\n"
        f"{_style_guidance(style)}\n\n"
        f"{json.dumps(payload, indent=2)}"
    )
    text = llm_complete(system, user, temperature=0.25, task="compose_slides")
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    parsed = _parse_scenes_json(text)
    if len(parsed) < 1:
        return scenes
    for idx, scene in enumerate(parsed):
        if idx < len(scenes):
            scene.chart_path = scenes[idx].chart_path
            scene.slide_image = scenes[idx].slide_image
            if scene.duration <= 0:
                scene.duration = scenes[idx].duration
    return _enrich_slide_scenes(parsed)


def _template_slides_script(topic: str, *, style: str = "executive") -> list[Scene]:
    label = topic_label(topic)
    style = normalize_slide_style(style)
    if style == "pitch":
        scenes = [
            Scene(
                title=f"{label} will reshape its market",
                narration=f"Open with the big opportunity around {label} and why now is the moment to act.",
                body="",
                captions=["Urgent market shift", "Window is open now"],
                image_query="startup pitch audience",
            ),
            Scene(
                title="Customers face a costly, painful problem",
                narration="Name the pain clearly — time lost, money wasted, or risk accepted today.",
                body="",
                captions=["Pain is widespread", "Status quo is expensive"],
                image_query="frustrated business team",
            ),
            Scene(
                title=f"Our approach to {label} is 10x better",
                narration="Explain the solution in one crisp sentence and what makes it defensible.",
                body="",
                captions=["Clear product wedge", "Hard to copy advantage"],
                image_query="product demo laptop",
            ),
            Scene(
                title="A large market is ready to buy",
                narration="Size the opportunity with a credible market framing and target segment.",
                body="",
                captions=["Growing demand", "Clear buyer persona"],
                image_query="market growth chart",
            ),
            Scene(
                title="Early traction proves the model",
                narration="Share proof points: pilots, revenue, engagement, or design partners.",
                body="",
                captions=["Validated with users", "Metrics trending up"],
                image_query="team celebration office",
            ),
            Scene(
                title="The ask: partner with us to scale",
                narration="Close with a specific call to action — funding, pilot, or next meeting.",
                body="",
                captions=["Join the next phase", "Let's talk this week"],
                image_query="handshake business deal",
            ),
        ]
    elif style == "academic":
        scenes = [
            Scene(
                title=f"Why {label} matters today",
                narration=f"Introduce {label}, its relevance, and the question this talk answers.",
                body="",
                captions=["Research context", "Talk objective"],
                image_query="university lecture hall",
            ),
            Scene(
                title="Background and prior work",
                narration="Summarize established knowledge and the gap your narrative addresses.",
                body="",
                captions=["Key prior findings", "Open research gap"],
                image_query="library research books",
            ),
            Scene(
                title=f"Core concepts in {label}",
                narration="Define the essential ideas the audience needs for the rest of the deck.",
                body="",
                captions=["Definitions", "Framework"],
                image_query="whiteboard equations",
            ),
            Scene(
                title="Evidence and examples",
                narration="Present supporting data, cases, or results that substantiate the argument.",
                body="",
                captions=["Data-backed claims", "Illustrative cases"],
                image_query="scientific data chart",
            ),
            Scene(
                title="Implications and limitations",
                narration="Discuss what the findings mean, caveats, and open questions.",
                body="",
                captions=["Practical implications", "Known limits"],
                image_query="panel discussion academics",
            ),
            Scene(
                title="Conclusion and future work",
                narration="Restate the main takeaway and suggest next research or application steps.",
                body="",
                captions=["Main takeaway", "Future directions"],
                image_query="graduation academic audience",
            ),
        ]
    else:
        scenes = [
            Scene(
                title=f"{label} is a strategic priority now",
                narration=f"Hook the room: why {label} deserves executive attention this quarter.",
                body="",
                captions=["High stakes decision", "Momentum is building"],
                image_query="executive boardroom",
            ),
            Scene(
                title="The core challenge we must solve",
                narration="Frame the problem in business terms — cost, risk, speed, or customer impact.",
                body="",
                captions=["Problem is measurable", "Delay increases risk"],
                image_query="business strategy meeting",
            ),
            Scene(
                title=f"Key insight: what changes with {label}",
                narration="Deliver the central insight or finding that reframes the conversation.",
                body="",
                captions=["Clear point of view", "Supported by evidence"],
                image_query="data dashboard office",
            ),
            Scene(
                title="What this means for our organization",
                narration="Translate the insight into operating impact, tradeoffs, and stakeholders affected.",
                body="",
                captions=["Operational impact", "Stakeholders affected"],
                image_query="team planning session",
            ),
            Scene(
                title="Recommended path forward",
                narration="Propose a focused plan with priorities, owners, and a realistic timeline.",
                body="",
                captions=["3 priorities", "90-day horizon"],
                image_query="project roadmap wall",
            ),
            Scene(
                title="Decision and next step",
                narration="End with a explicit ask — approve, fund, pilot, or schedule a follow-up.",
                body="",
                captions=["Decision needed", "Next meeting scheduled"],
                image_query="handshake executives",
            ),
        ]
    return _enrich_slide_scenes(scenes)


def _llm_slides_enrich_image_keywords(topic: str, scenes: list[Scene]) -> list[Scene]:
    missing = [scene for scene in scenes if not scene.image_keywords]
    if not missing:
        return scenes
    try:
        from arka.llm.fallback import llm_complete
    except ImportError:
        return scenes

    from arka.media.compose_video import compact_photo_query

    payload = [
        {
            "title": scene.title,
            "narration": scene.narration[:300],
            "image_query": scene.image_query,
        }
        for scene in missing
    ]
    system = (
        "You pick stock-photo search keywords for presentation slides. "
        "Return ONLY JSON array: "
        '[{"title":"same as input", "image_keywords":["conference room", "team whiteboard"]}]. '
        "Provide 3-5 image_keywords per slide. Each keyword must be 2-3 visual nouns (no sentences)."
    )
    user = f"Topic: {topic_label(topic)}\n\n{json.dumps(payload, indent=2)}"
    text = llm_complete(system, user, temperature=0.2, task="compose_slides")
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    try:
        rows = json.loads(text)
    except json.JSONDecodeError:
        return scenes
    by_title = {str(row.get("title") or "").strip(): row for row in rows if isinstance(row, dict)}
    for scene in scenes:
        row = by_title.get(scene.title.strip())
        if not row:
            continue
        raw = row.get("image_keywords") or row.get("keywords") or row.get("images")
        if isinstance(raw, list):
            scene.image_keywords = [
                compact_photo_query(str(item).strip())
                for item in raw
                if str(item).strip()
            ]
            scene.image_keywords = [kw for kw in scene.image_keywords if kw]
        if not scene.image_query.strip():
            query = row.get("image_query") or row.get("image_search")
            if isinstance(query, str) and query.strip():
                scene.image_query = compact_photo_query(query.strip())
    return scenes


def _default_format() -> str:
    raw = (_env("SLIDES_DEFAULT_FORMAT") or "pptx").lower()
    return normalize_format(raw) or "pptx"


def normalize_format(name: str) -> str:
    return FORMAT_ALIASES.get((name or "").strip().lower(), "")


def parse_formats_arg(value: str | None) -> list[str]:
    raw = (value or "auto").strip().lower()
    if raw in {"", "auto"}:
        return [_default_format()]
    if raw == "all":
        return list(SLIDE_FORMATS)
    formats: list[str] = []
    for token in re.split(r"[,+\s]+", raw):
        token = token.strip()
        if not token:
            continue
        fmt = normalize_format(token)
        if not fmt:
            known = ", ".join(SLIDE_FORMATS + ("all", "auto"))
            raise SystemExit(f"Unknown slide format {token!r}. Choose: {known}")
        if fmt not in formats:
            formats.append(fmt)
    return formats or [_default_format()]


def _output_paths(output: Path, formats: list[str]) -> dict[str, Path]:
    stem = output.with_suffix("")
    if len(formats) == 1:
        fmt = formats[0]
        return {fmt: stem.with_suffix(FORMAT_EXTENSIONS[fmt])}
    return {fmt: stem.with_suffix(FORMAT_EXTENSIONS[fmt]) for fmt in formats}


def _default_output(topic: str, fmt: str = "pptx") -> Path:
    slug = re.sub(r"[^a-z0-9]+", "-", topic_label(topic).lower())[:40].strip("-") or "slides"
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    env_dir = _env("SLIDES_OUTPUT_DIR") or _env("VIDEO_OUTPUT_DIR") or _env("IMAGE_OUTPUT_DIR")
    out_dir = Path(env_dir).expanduser() if env_dir else Path.home() / "Documents" / "arka-slides"
    out_dir.mkdir(parents=True, exist_ok=True)
    ext = FORMAT_EXTENSIONS.get(fmt, ".pptx")
    return out_dir / f"{slug}-{ts}{ext}"


def _strip_format_from_text(text: str) -> tuple[str, str | None]:
    """Remove format hints from NL; return (cleaned, format or None)."""
    t = text.strip()
    if not t:
        return "", None

    alias_pat = "|".join(re.escape(k) for k in sorted(FORMAT_ALIASES, key=len, reverse=True))

    def _alias(word: str) -> str:
        return FORMAT_ALIASES.get(word.lower(), "")

    found: str | None = None

    m = re.search(rf"(?i)\b(?:as|in|to|into|export(?:ed)?)\s+({alias_pat})\s*$", t)
    if m:
        found = _alias(m.group(1)) or found
        t = t[: m.start()].strip()

    m = re.search(rf"(?i)^({alias_pat})\s+(?=(?:slide|slides|presentation|deck|slideshow)\b)", t)
    if m:
        found = found or _alias(m.group(1))
        t = (t[: m.start()] + t[m.end() :]).strip()

    m = re.search(rf"(?i)\b({alias_pat})\s+(?=(?:slide|slides|presentation|deck|slideshow)\b)", t)
    if m:
        found = found or _alias(m.group(1))
        t = (t[: m.start()] + t[m.end() :]).strip()

    t = re.sub(r"\s{2,}", " ", t).strip(" ,;")
    return t, found


def _extract_format_from_nl(text: str) -> str | None:
    _, fmt = _strip_format_from_text(text)
    return fmt


def _require_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as exc:
        raise SystemExit(
            "python-pptx is required for slide decks.\n"
            "Install: pip install 'arka-agent[video]'  or  pip install python-pptx Pillow"
        ) from exc
    return Presentation, Emu


class SlideExportError(Exception):
    """Raised when a slide export file fails validation."""

    def __init__(self, fmt: str, path: Path, reason: str) -> None:
        self.fmt = fmt
        self.path = path
        self.reason = reason
        super().__init__(f"{fmt} export invalid ({path.name}): {reason}")


@dataclass
class ExportBatch:
    saved: list[Path]
    outputs: dict[str, str]
    failed: dict[str, str] = field(default_factory=dict)
    fallback_md: Path | None = None


def _sanitize_ooxml_text(text: str) -> str:
    """Strip characters illegal in XML 1.0 text (can corrupt OOXML inside pptx)."""
    return "".join(
        ch
        for ch in text
        if ch in "\t\n\r"
        or (0x20 <= ord(ch) <= 0xD7FF)
        or (0xE000 <= ord(ch) <= 0xFFFD)
    )


def _pptx_include_notes() -> bool:
    raw = (_env("SLIDES_PPTX_NOTES") or "").lower()
    return raw in {"1", "true", "yes", "on"}


EMU_PER_INCH = 914400
# Keynote imports standard PowerPoint slide sizes; avoid mapping video pixels 1:1 to EMU.
_PPTX_WIDESCREEN_WIDTH_IN = 13.333
_PPTX_WIDESCREEN_HEIGHT_IN = 7.5
_PPTX_MAX_SLIDE_INCHES = 14.0


def _pptx_slide_dimensions(_cfg: VideoConfig | None = None) -> tuple[int, int]:
    """Return (width_emu, height_emu) using standard 16:9 widescreen PowerPoint size."""
    _require_pptx()
    from pptx.util import Inches

    return int(Inches(_PPTX_WIDESCREEN_WIDTH_IN)), int(Inches(_PPTX_WIDESCREEN_HEIGHT_IN))


def _prepare_pptx_image_stream(path: Path) -> io.BytesIO:
    """Re-encode a slide PNG as clean RGB for reliable OOXML / Keynote embedding."""
    from arka.media.compose_video import _require_pillow

    Image, *_ = _require_pillow()
    with Image.open(path) as img:
        if img.mode == "RGBA":
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[3])
            rgb = bg
        elif img.mode != "RGB":
            rgb = img.convert("RGB")
        else:
            rgb = img.copy()
        buf = io.BytesIO()
        rgb.save(buf, format="PNG")
        buf.seek(0)
        return buf


def _validate_slide_image(path: Path) -> None:
    if not path.is_file():
        raise ValueError(f"slide image missing: {path}")
    size = path.stat().st_size
    if size < 64:
        raise ValueError(f"slide image too small ({size} bytes): {path}")
    from arka.media.compose_video import _require_pillow

    Image, *_ = _require_pillow()
    with Image.open(path) as img:
        img.verify()
    with Image.open(path) as img:
        if img.width < 8 or img.height < 8:
            raise ValueError(f"slide image dimensions invalid: {img.width}x{img.height}")


def _validate_pptx_media_bytes(name: str, data: bytes) -> None:
    if len(data) < 24:
        raise ValueError(f"media too small: {name}")
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        if b"IEND" not in data[-32:]:
            raise ValueError(f"truncated PNG: {name}")
        return
    if data[:2] == b"\xff\xd8":
        if b"\xff\xd9" not in data[-4:]:
            raise ValueError(f"truncated JPEG: {name}")
        return
    raise ValueError(f"unsupported media type in {name}")


def _validate_pptx_file(path: Path) -> None:
    data = path.read_bytes()
    if len(data) < 512:
        raise ValueError("file too small to be a pptx")
    if not data.startswith(b"PK\x03\x04"):
        raise ValueError("not a ZIP archive (expected OOXML pptx container)")
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        if "[Content_Types].xml" not in names:
            raise ValueError("missing [Content_Types].xml")
        if "ppt/presentation.xml" not in names:
            raise ValueError("missing ppt/presentation.xml")
        if "_rels/.rels" not in names:
            raise ValueError("missing package _rels/.rels")
        if not any(n.startswith("ppt/slides/slide") and n.endswith(".xml") for n in names):
            raise ValueError("missing ppt/slides/slide*.xml")
        bad = zf.testzip()
        if bad:
            raise ValueError(f"corrupt zip entry: {bad}")
        for name in names:
            if not name.endswith(".xml"):
                continue
            raw = zf.read(name)
            text = raw.decode("utf-8")
            illegal = [c for c in text if ord(c) < 0x20 and c not in "\t\n\r"]
            if illegal:
                raise ValueError(f"illegal XML characters in {name}")
            ET.fromstring(raw)
        for name in names:
            if name.startswith("ppt/media/"):
                _validate_pptx_media_bytes(name, zf.read(name))
    Presentation, _ = _require_pptx()
    prs = Presentation(str(path))
    if len(prs.slides) < 1:
        raise ValueError("pptx contains no slides")
    w_in = prs.slide_width / EMU_PER_INCH
    h_in = prs.slide_height / EMU_PER_INCH
    if w_in > _PPTX_MAX_SLIDE_INCHES or h_in > _PPTX_MAX_SLIDE_INCHES:
        raise ValueError(
            f"slide size {w_in:.2f}x{h_in:.2f} in exceeds Keynote-safe maximum "
            f"({_PPTX_MAX_SLIDE_INCHES} in)"
        )
    if w_in < 5.0 or h_in < 5.0:
        raise ValueError(f"slide size {w_in:.2f}x{h_in:.2f} in is too small")
    aspect = w_in / h_in
    if not (1.2 <= aspect <= 1.9):
        raise ValueError(f"non-standard slide aspect ratio {aspect:.2f}")


def _validate_pdf_file(path: Path) -> None:
    if not path.read_bytes()[:5].startswith(b"%PDF-"):
        raise ValueError("invalid PDF header")


def _validate_html_file(path: Path) -> None:
    head = path.read_text(encoding="utf-8", errors="replace")[:4096].lower()
    if "<html" not in head and "<!doctype html" not in head:
        raise ValueError("missing HTML document structure")


def _validate_markdown_file(path: Path) -> None:
    text = path.read_text(encoding="utf-8")
    if not text.strip():
        raise ValueError("empty markdown export")
    if "# " not in text and not text.lstrip().startswith("---"):
        raise ValueError("markdown export has no slide headings")


def _validate_json_file(path: Path) -> None:
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise ValueError("JSON export must be an object")
    if not data.get("scenes"):
        raise ValueError("JSON export missing scenes")


def _export_validator(fmt: str):
    return {
        "pptx": _validate_pptx_file,
        "pdf": _validate_pdf_file,
        "html": _validate_html_file,
        "md": _validate_markdown_file,
        "json": _validate_json_file,
    }.get(fmt)


def _remove_invalid_export(path: Path) -> None:
    try:
        if path.is_file():
            path.unlink()
    except OSError:
        pass


def _commit_export(path: Path, fmt: str) -> Path:
    validator = _export_validator(fmt)
    if validator:
        try:
            validator(path)
        except Exception as exc:
            _remove_invalid_export(path)
            raise SlideExportError(fmt, path, str(exc)) from exc
    return path


def _atomic_export(build_fn, path: Path, fmt: str) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    partial = path.with_name(f"{path.name}.partial")
    _remove_invalid_export(partial)
    try:
        build_fn(partial)
        _commit_export(partial, fmt)
        partial.replace(path)
        return path
    except Exception:
        _remove_invalid_export(partial)
        _remove_invalid_export(path)
        raise


def extract_slides_topic(text: str) -> str:
    """Pull the subject out of NL like 'make slides about kubernetes'."""
    t = text.strip().strip("'\"")
    if not t:
        return ""
    t, _ = _strip_format_from_text(t)
    t = _strip_style_from_text(t)
    patterns = [
        r"(?i)(?:^|\b)(?:make|create|compose|build|render|produce|generate|arka)\s+"
        r"(?:a\s+|an\s+)?(?:\d+\s+)?(?:slide|slides|presentation|deck|powerpoint|pptx?|pdf|html|markdown|marp)\s+"
        r"(?:on|about|for|explaining|covering)\s+(.+)$",
        r"(?i)(?:^|\b)(?:slide|slides|presentation|deck)\s+(?:on|about|for|explaining|covering)\s+(.+)$",
        r"(?i)^compose\s+(?:a\s+)?(?:slide|slides|presentation|deck)\s+"
        r"(?:on|about|for|explaining|covering)\s+(.+)$",
        r"(?i)(?:^|\b)(?:pdf|html|markdown|marp)\s+(?:slide|slides|presentation|deck)\s+"
        r"(?:on|about|for|explaining|covering)\s+(.+)$",
        r"(?i)(?:^|\b)(?:executive|academic|pitch)\s+"
        r"(?:slide|slides|presentation|deck)\s+(?:on|about|for|explaining|covering)\s+(.+)$",
    ]
    for pat in patterns:
        m = re.search(pat, t)
        if m:
            topic = m.group(1).strip().strip("'\"")
            topic = re.sub(r"(?i)\s+(?:with\s+llm|please)$", "", topic).strip()
            topic, _ = _strip_format_from_text(topic)
            topic = _strip_style_from_text(topic)
            if topic:
                return topic
    cleaned, _ = _strip_format_from_text(t)
    return normalize_topic(_strip_style_from_text(cleaned))


def nl_to_argv_convert(text: str) -> list[str]:
    """Parse NL like 'convert slides.pptx to pdf' → convert argv."""
    t = text.strip().strip("'\"")
    if not t:
        return []

    alias_pat = "|".join(re.escape(k) for k in sorted(FORMAT_ALIASES, key=len, reverse=True))

    if not re.search(
        r"(?i)\b(?:convert|export|transform)\s+(?:my\s+)?(?:presentation|slides|deck|slideshow)\b",
        t,
    ) and not re.search(
        r"(?i)\bconvert\s+\S+\.(?:pptx|pdf|html|md|json|markdown|marp|ppt|htm)\b",
        t,
    ):
        if not re.search(
            r"(?i)\bconvert\s+['\"]?.+\.(?:pptx|pdf|html|md|json|markdown|marp|ppt|htm)",
            t,
        ):
            return []

    fmt: str | None = None
    for pat in (
        rf"(?i)(?:-+\s*)?(?:format|output)\s+(?:is\s+)?({alias_pat}|all)\b",
        rf"(?i)\b(?:to|into|as)\s+(?:a\s+)?({alias_pat}|all)\b",
    ):
        m = re.search(pat, t)
        if m:
            token = m.group(1).lower()
            fmt = "all" if token == "all" else normalize_format(token)
            t = (t[: m.start()] + t[m.end() :]).strip()

    input_path = ""
    for pat in (
        r"(?i)\bconvert\s+(.+?)\s+(?:to|into|as)\s+",
        r"(?i)\b(?:convert|export|transform)\s+(?:my\s+)?(?:presentation|slides|deck|slideshow)\s+(?:to|into|as)\s+",
        r"(?i)\bconvert\s+(.+)$",
        r"(?i)\b(?:export|transform)\s+(.+?)\s+(?:to|into|as)\s+",
    ):
        m = re.search(pat, t)
        if m:
            candidate = m.group(1).strip().strip("'\"")
            if candidate and not normalize_format(candidate):
                input_path = candidate
                break

    if not input_path:
        file_m = re.search(
            r"(?i)(['\"]?[\w./~-]+\.(?:pptx|pdf|html|htm|md|markdown|marp|json|ppt)['\"]?)",
            t,
        )
        if file_m:
            input_path = file_m.group(1).strip("'\"")

    if not input_path:
        return []

    argv = ["convert", input_path]
    if fmt:
        argv.extend(["--format", fmt])
    return argv


def nl_to_argv(text: str) -> list[str]:
    t = text.strip()
    if not t:
        return []

    convert_argv = nl_to_argv_convert(t)
    if convert_argv:
        return convert_argv

    slide_intent = re.search(
        r"(?i)(?:^|\b)(?:make|create|compose|build|render|produce|generate|arka)\s+"
        r"(?:a\s+|an\s+)?(?:\d+\s+)?(?:slide|slides|presentation|deck|powerpoint|pptx?|pdf|html|markdown|marp)\s+"
        r"(?:on|about|for|explaining|covering)\s+\S",
        t,
    ) or re.search(
        r"(?i)(?:^|\b)(?:slide|slides|presentation|deck)\s+(?:on|about|for|explaining|covering)\s+\S",
        t,
    ) or re.search(
        r"(?i)^compose\s+(?:a\s+)?(?:slide|slides|presentation|deck)\s+"
        r"(?:on|about|for|explaining|covering)\s+\S",
        t,
    ) or re.search(
        r"(?i)(?:^|\b)(?:pdf|html|markdown|marp)\s+(?:slide|slides|presentation|deck)\s+"
        r"(?:on|about|for|explaining|covering)\s+\S",
        t,
    ) or re.search(
        r"(?i)(?:^|\b)(?:executive|academic|pitch)\s+"
        r"(?:slide|slides|presentation|deck)\s+(?:on|about|for|explaining|covering)\s+\S",
        t,
    )
    if not slide_intent:
        return []

    topic = extract_slides_topic(t)
    if not topic:
        return []
    argv = ["compose", "--topic", topic]
    fmt = _extract_format_from_nl(t)
    if fmt:
        argv.extend(["--format", fmt])
    style = extract_slide_style(t)
    if re.search(r"(?i)\b(?:executive|academic|pitch)\b", t):
        argv.extend(["--style", style])
    if re.search(r"(?i)\b(llm|write script)\b", t):
        argv.append("--llm")
    return argv


def _render_scene_png(
    scene: Scene,
    *,
    topic: str,
    work: Path,
    index: int,
    cfg: VideoConfig,
    used_photo_ids: set[str],
    credits: list[dict],
) -> Path:
    from arka.media.chart_slide import render_scene_visual, render_title_slide, scene_has_chart_visual

    out = work / f"slide-{index:02d}.png"
    if scene_has_chart_visual(scene):
        return render_scene_visual(scene, work, cfg, index=index)

    if any_source_available():
        query = _scene_search_query(scene, topic)
        try:
            photo = _fetch_scene_photo(
                query,
                used_photo_ids,
                orientation=cfg.orientation,
                segment_idx=0,
            )
            scene.photo = photo
            credits.append(
                {
                    "title": scene.title,
                    "query": query,
                    "source": getattr(photo, "source", ""),
                    "author": getattr(photo, "author", ""),
                    "url": getattr(photo, "url", ""),
                }
            )
            _render_photo_slide(
                scene,
                photo,
                work=work,
                scene_idx=index,
                segment_idx=0,
                slide=out,
                cfg=cfg,
            )
            return out
        except SystemExit:
            pass

    return render_title_slide(scene, out, cfg)


def _build_pptx(
    slide_images: list[tuple[Path, Scene]],
    output: Path,
    *,
    topic: str,
    cfg: VideoConfig,
) -> Path:
    Presentation, _ = _require_pptx()
    include_notes = _pptx_include_notes()
    for png_path, _ in slide_images:
        _validate_slide_image(png_path)

    slide_w, slide_h = _pptx_slide_dimensions(cfg)

    def _write(dest: Path) -> None:
        prs = Presentation()
        prs.slide_width = slide_w
        prs.slide_height = slide_h
        blank_layout = prs.slide_layouts[6]

        for png_path, scene in slide_images:
            slide = prs.slides.add_slide(blank_layout)
            try:
                stream = _prepare_pptx_image_stream(png_path)
                slide.shapes.add_picture(
                    stream,
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
            except Exception as exc:
                print(
                    f"  Skipping slide image ({png_path.name}): {exc}",
                    file=sys.stderr,
                )
            if include_notes:
                notes = _sanitize_ooxml_text((scene.narration or scene.body or "").strip())
                if notes:
                    slide.notes_slide.notes_text_frame.text = notes

        prs.save(str(dest))

    return _atomic_export(_write, output, "pptx")


def _build_pdf(
    slide_images: list[tuple[Path, Scene]],
    output: Path,
    **_,
) -> Path:
    from arka.media.compose_video import _require_pillow

    Image, *_ = _require_pillow()
    images = []
    for png_path, _ in slide_images:
        _validate_slide_image(png_path)
        img = Image.open(png_path)
        if img.mode != "RGB":
            img = img.convert("RGB")
        images.append(img)
    if not images:
        raise SystemExit("No slide images to export as PDF.")

    def _write(dest: Path) -> None:
        images[0].save(
            str(dest),
            "PDF",
            save_all=True,
            append_images=images[1:],
            resolution=100.0,
        )

    return _atomic_export(_write, output, "pdf")


def _png_data_uri(path: Path) -> str:
    import base64

    data = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:image/png;base64,{data}"


def _build_html(
    slide_images: list[tuple[Path, Scene]],
    output: Path,
    *,
    topic: str,
    **_,
) -> Path:
    slides_html = []
    for i, (png_path, scene) in enumerate(slide_images):
        active = " active" if i == 0 else ""
        notes = (scene.narration or scene.body or "").strip()
        note_html = f'<p class="notes">{_html_escape(notes)}</p>' if notes else ""
        slides_html.append(
            f'<section class="slide{active}" data-index="{i}">'
            f'<img src="{_png_data_uri(png_path)}" alt="{_html_escape(scene.title)}">'
            f"{note_html}</section>"
        )
    title = _html_escape(topic_label(topic))
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title}</title>
<style>
  * {{ box-sizing: border-box; }}
  body {{ margin: 0; background: #111; color: #eee; font-family: system-ui, sans-serif; }}
  .deck {{ width: 100vw; height: 100vh; position: relative; overflow: hidden; }}
  .slide {{ display: none; width: 100%; height: 100%; align-items: center; justify-content: center; flex-direction: column; }}
  .slide.active {{ display: flex; }}
  .slide img {{ max-width: 100%; max-height: calc(100% - 3rem); object-fit: contain; }}
  .notes {{ position: absolute; bottom: 0; left: 0; right: 0; margin: 0; padding: 0.75rem 1rem; background: rgba(0,0,0,0.65); font-size: 0.9rem; }}
  .hint {{ position: fixed; top: 0.5rem; right: 0.5rem; opacity: 0.5; font-size: 0.75rem; }}
</style>
</head>
<body>
<div class="deck" id="deck">
{"".join(slides_html)}
</div>
<p class="hint">← → or Space</p>
<script>
(function() {{
  const slides = Array.from(document.querySelectorAll('.slide'));
  let idx = 0;
  function show(i) {{
    idx = (i + slides.length) % slides.length;
    slides.forEach((s, n) => s.classList.toggle('active', n === idx));
  }}
  document.addEventListener('keydown', (e) => {{
    if (['ArrowRight', 'ArrowDown', ' ', 'PageDown'].includes(e.key)) {{ e.preventDefault(); show(idx + 1); }}
    if (['ArrowLeft', 'ArrowUp', 'PageUp'].includes(e.key)) {{ e.preventDefault(); show(idx - 1); }}
  }});
  document.getElementById('deck').addEventListener('click', () => show(idx + 1));
}})();
</script>
</body>
</html>
"""
    output.parent.mkdir(parents=True, exist_ok=True)
    partial = output.with_name(f"{output.name}.partial")
    _remove_invalid_export(partial)
    try:
        partial.write_text(html, encoding="utf-8")
        _commit_export(partial, "html")
        partial.replace(output)
    except Exception:
        _remove_invalid_export(partial)
        _remove_invalid_export(output)
        raise
    return output


def _html_escape(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _build_markdown(
    slide_images: list[tuple[Path, Scene]],
    output: Path,
    *,
    topic: str,
    scenes: list[Scene],
    **_,
) -> Path:
    _ = slide_images
    lines = [
        "---",
        f"title: {topic_label(topic)}",
        "marp: true",
        "---",
        "",
    ]
    for i, scene in enumerate(scenes):
        if i > 0:
            lines.extend(["", "---", ""])
        lines.append(f"# {scene.title}")
        body = (scene.body or "").strip()
        if body:
            lines.extend(["", body])
        captions = [c.strip() for c in (scene.captions or []) if c and c.strip()]
        if captions:
            lines.append("")
            lines.extend(f"- {c}" for c in captions)
        notes = (scene.narration or "").strip()
        if notes:
            lines.extend(["", f"<!-- speaker notes: {notes} -->"])
    output.parent.mkdir(parents=True, exist_ok=True)
    partial = output.with_name(f"{output.name}.partial")
    _remove_invalid_export(partial)
    try:
        partial.write_text("\n".join(lines) + "\n", encoding="utf-8")
        _commit_export(partial, "md")
        partial.replace(output)
    except Exception:
        _remove_invalid_export(partial)
        _remove_invalid_export(output)
        raise
    return output


def _metadata_payload(
    *,
    topic: str,
    scenes: list[Scene],
    credits: list[dict],
    outputs: dict[str, str],
) -> dict:
    return {
        "topic": topic,
        "scenes": [
            {
                "title": s.title,
                "narration": s.narration,
                "body": s.body,
                "captions": s.captions,
                "image_query": s.image_query,
                "image_keywords": s.image_keywords,
                "chart": s.chart,
                "chart_path": s.chart_path,
                "slide_image": s.slide_image,
            }
            for s in scenes
        ],
        "photo_credits": credits,
        "outputs": outputs,
        "formats": list(outputs),
        "source": "arka-compose-slides",
    }


def _build_json_export(
    slide_images: list[tuple[Path, Scene]],
    output: Path,
    *,
    topic: str,
    scenes: list[Scene],
    credits: list[dict],
    outputs: dict[str, str] | None = None,
    **_,
) -> Path:
    _ = slide_images
    payload = _metadata_payload(
        topic=topic,
        scenes=scenes,
        credits=credits,
        outputs=outputs or {output.suffix.lstrip("."): str(output)},
    )
    output.parent.mkdir(parents=True, exist_ok=True)
    partial = output.with_name(f"{output.name}.partial")
    _remove_invalid_export(partial)
    try:
        partial.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")
        _commit_export(partial, "json")
        partial.replace(output)
    except Exception:
        _remove_invalid_export(partial)
        _remove_invalid_export(output)
        raise
    return output


_EXPORTERS = {
    "pptx": _build_pptx,
    "pdf": _build_pdf,
    "html": _build_html,
    "md": _build_markdown,
    "json": _build_json_export,
}


def _export_formats(
    slide_images: list[tuple[Path, Scene]],
    output_paths: dict[str, Path],
    *,
    topic: str,
    scenes: list[Scene],
    cfg: VideoConfig,
    credits: list[dict],
) -> ExportBatch:
    saved: list[Path] = []
    errors: list[str] = []
    outputs_map: dict[str, str] = {}
    failed: dict[str, str] = {}
    fallback_md: Path | None = None

    def _run_export(fmt: str, path: Path) -> Path:
        exporter = _EXPORTERS[fmt]
        if fmt in {"pptx", "pdf", "html"}:
            return exporter(slide_images, path, topic=topic, cfg=cfg)
        return exporter(
            slide_images,
            path,
            topic=topic,
            scenes=scenes,
            cfg=cfg,
        )

    def _try_markdown_fallback(pptx_path: Path) -> Path | None:
        if "md" in output_paths:
            return None
        md_path = pptx_path.with_suffix(".md")
        try:
            saved_path = _build_markdown(
                slide_images,
                md_path,
                topic=topic,
                scenes=scenes,
                cfg=cfg,
            )
            saved.append(saved_path)
            outputs_map["md"] = str(saved_path)
            return saved_path
        except Exception as exc:
            _remove_invalid_export(md_path)
            errors.append(f"md fallback: {exc}")
            return None

    for fmt, path in output_paths.items():
        if fmt == "json":
            continue
        try:
            saved_path = _run_export(fmt, path)
            saved.append(saved_path)
            outputs_map[fmt] = str(saved_path)
        except SystemExit as exc:
            msg = str(exc).strip() or f"{fmt} export unavailable"
            failed[fmt] = msg
            _remove_invalid_export(path)
            if fmt == "pptx":
                fallback_md = _try_markdown_fallback(path) or fallback_md
            if len(output_paths) == 1 and not fallback_md:
                raise
            errors.append(f"{fmt}: {msg}")
            print(f"  Skipping {fmt} export — {msg}", file=sys.stderr)
        except SlideExportError as exc:
            failed[fmt] = exc.reason
            if fmt == "pptx":
                fallback_md = _try_markdown_fallback(path) or fallback_md
            if len(output_paths) == 1 and not fallback_md:
                raise SystemExit(str(exc)) from exc
            errors.append(f"{fmt}: {exc.reason}")
            print(f"  Skipping {fmt} export — {exc.reason}", file=sys.stderr)
        except Exception as exc:
            failed[fmt] = str(exc)
            _remove_invalid_export(path)
            if fmt == "pptx":
                fallback_md = _try_markdown_fallback(path) or fallback_md
            if len(output_paths) == 1 and not fallback_md:
                raise
            errors.append(f"{fmt}: {exc}")
            print(f"  Skipping {fmt} export — {exc}", file=sys.stderr)

    if not saved and "json" not in output_paths:
        detail = "; ".join(errors) if errors else "no exporters succeeded"
        raise SystemExit(f"No slide formats exported ({detail}).")

    if "json" in output_paths:
        try:
            saved_path = _build_json_export(
                slide_images,
                output_paths["json"],
                topic=topic,
                scenes=scenes,
                credits=credits,
                outputs=outputs_map,
                cfg=cfg,
            )
            saved.append(saved_path)
            outputs_map["json"] = str(saved_path)
        except (SystemExit, SlideExportError) as exc:
            msg = str(exc).strip() or "json export unavailable"
            failed["json"] = msg
            _remove_invalid_export(output_paths["json"])
            if len(output_paths) == 1 and not saved:
                raise SystemExit(f"No slide formats exported ({msg}).") from exc
            errors.append(f"json: {msg}")
            print(f"  Skipping json export — {msg}", file=sys.stderr)
        except Exception as exc:
            failed["json"] = str(exc)
            _remove_invalid_export(output_paths["json"])
            if len(output_paths) == 1 and not saved:
                raise
            errors.append(f"json: {exc}")
            print(f"  Skipping json export — {exc}", file=sys.stderr)

    if not saved:
        detail = "; ".join(errors) if errors else "no exporters succeeded"
        raise SystemExit(f"No slide formats exported ({detail}).")

    return ExportBatch(saved=saved, outputs=outputs_map, failed=failed, fallback_md=fallback_md)


def compose(
    scenes: list[Scene],
    *,
    output: Path,
    topic: str,
    cfg: VideoConfig | None = None,
    formats: list[str] | None = None,
) -> ExportBatch:
    if not scenes:
        raise SystemExit("No scenes to render.")
    cfg = cfg or load_config()
    formats = formats or [_default_format()]
    from arka.media.chart_slide import scene_has_chart_visual

    needs_photos = any(not scene_has_chart_visual(s) for s in scenes)
    if needs_photos:
        _attach_photo_queries(scenes, topic)
        if not any_source_available():
            print(
                "  No stock photo API keys — using title slides only.",
                file=sys.stderr,
            )

    work = Path(tempfile.mkdtemp(prefix="arka-slides-"))
    used_photo_ids: set[str] = set()
    credits: list[dict] = []
    slide_images: list[tuple[Path, Scene]] = []

    try:
        for i, scene in enumerate(scenes):
            print(f"  Slide {i + 1}/{len(scenes)}: {scene.title}", file=sys.stderr)
            png = _render_scene_png(
                scene,
                topic=topic,
                work=work,
                index=i,
                cfg=cfg,
                used_photo_ids=used_photo_ids,
                credits=credits,
            )
            _validate_slide_image(png)
            slide_images.append((png, scene))

        output_paths = _output_paths(output, formats)
        batch = _export_formats(
            slide_images,
            output_paths,
            topic=topic,
            scenes=scenes,
            cfg=cfg,
            credits=credits,
        )

        if "json" not in batch.outputs:
            sidecar = output.with_suffix(".meta.json")
            sidecar.write_text(
                json.dumps(
                    _metadata_payload(
                        topic=topic,
                        scenes=scenes,
                        credits=credits,
                        outputs=batch.outputs,
                    ),
                    indent=2,
                )
                + "\n",
                encoding="utf-8",
            )
        return batch
    finally:
        import shutil

        shutil.rmtree(work, ignore_errors=True)


@dataclass
class LoadedDeck:
    topic: str
    scenes: list[Scene]
    slide_images: list[tuple[Path, Scene]] | None = None
    credits: list[dict] = field(default_factory=list)


def detect_format(path: Path, from_fmt: str | None = None) -> str:
    if from_fmt:
        fmt = normalize_format(from_fmt)
        if not fmt:
            known = ", ".join(SLIDE_FORMATS)
            raise SystemExit(f"Unknown --from format {from_fmt!r}. Choose: {known}")
        return fmt
    ext = path.suffix.lstrip(".").lower()
    fmt = normalize_format(ext)
    if fmt:
        return fmt
    known = ", ".join(SLIDE_FORMATS)
    raise SystemExit(
        f"Cannot detect slide format from {path.name!r} (extension {path.suffix!r}). "
        f"Use --from. Supported: {known}"
    )


def _find_sidecar(path: Path) -> Path | None:
    for candidate in (
        path.with_suffix(".meta.json"),
        path.with_suffix(".json"),
        path.parent / f"{path.stem}.meta.json",
    ):
        if candidate.is_file() and candidate.resolve() != path.resolve():
            try:
                data = json.loads(candidate.read_text(encoding="utf-8"))
            except json.JSONDecodeError:
                continue
            if isinstance(data, dict) and data.get("scenes"):
                return candidate
    return None


def _decode_data_uri(data_uri: str) -> bytes:
    if "," not in data_uri:
        raise ValueError("invalid data URI")
    _, encoded = data_uri.split(",", 1)
    return base64.b64decode(encoded)


def _decode_data_uri_to_file(data_uri: str, dest: Path) -> Path:
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(_decode_data_uri(data_uri))
    return dest


def _parse_markdown_deck(text: str) -> tuple[str, list[Scene]]:
    topic = ""
    body = text
    if body.startswith("---"):
        parts = body.split("---", 2)
        if len(parts) >= 3:
            for line in parts[1].splitlines():
                if line.strip().lower().startswith("title:"):
                    topic = line.split(":", 1)[1].strip().strip("'\"")
            body = parts[2]

    chunks = [chunk.strip() for chunk in re.split(r"\n---\n", body.strip()) if chunk.strip()]
    scenes: list[Scene] = []
    for chunk in chunks:
        if chunk.startswith("---"):
            continue
        title = ""
        body_lines: list[str] = []
        captions: list[str] = []
        notes = ""
        for line in chunk.splitlines():
            stripped = line.strip()
            if re.match(r"^#{1,6}\s", stripped):
                title = re.sub(r"^#{1,6}\s*", "", stripped).strip()
            elif stripped.startswith("- "):
                captions.append(stripped[2:].strip())
            elif stripped.lower().startswith("<!-- speaker notes:"):
                notes = stripped.split(":", 1)[1].strip().rstrip("-->").strip()
            elif not stripped.startswith("<!--"):
                body_lines.append(line)
        body_text = "\n".join(body_lines).strip()
        if not title and not body_text and not captions:
            continue
        scenes.append(
            Scene(
                title=title or f"Slide {len(scenes) + 1}",
                body=body_text,
                narration=notes,
                captions=captions,
            )
        )
    return topic, scenes


def _load_json_deck(path: Path) -> LoadedDeck:
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise SystemExit(f"Expected JSON object in {path.name}")
    topic = str(data.get("topic") or "").strip()
    scenes = _enrich_scenes(_parse_scenes_json(json.dumps(data)))
    if not scenes:
        raise SystemExit(f"No scenes found in {path.name}")
    credits = data.get("photo_credits") if isinstance(data.get("photo_credits"), list) else []
    return LoadedDeck(
        topic=topic or scenes[0].title,
        scenes=scenes,
        credits=credits,
    )


def _load_markdown_deck(path: Path) -> LoadedDeck:
    topic, scenes = _parse_markdown_deck(path.read_text(encoding="utf-8"))
    if not scenes:
        raise SystemExit(f"No slides found in markdown: {path.name}")
    return LoadedDeck(topic=topic or path.stem, scenes=scenes)


def _load_html_deck(path: Path, work: Path) -> LoadedDeck:
    sidecar = _find_sidecar(path)
    if sidecar:
        return _load_json_deck(sidecar)

    html = path.read_text(encoding="utf-8")
    title_m = re.search(r"<title>([^<]+)</title>", html, re.I)
    topic = html_module.unescape(title_m.group(1).strip()) if title_m else path.stem

    scenes: list[Scene] = []
    slide_images: list[tuple[Path, Scene]] = []
    section_re = re.compile(
        r'<section[^>]*class="[^"]*slide[^"]*"[^>]*data-index="(\d+)"[^>]*>(.*?)</section>',
        re.S | re.I,
    )
    matches = list(section_re.finditer(html))
    if not matches:
        raise SystemExit(
            f"Could not parse slides from HTML: {path.name}\n"
            "Expected arka HTML sections, or place a .meta.json sidecar beside the file."
        )

    for match in matches:
        idx = int(match.group(1))
        body = match.group(2)
        img_m = re.search(r'<img[^>]+alt="([^"]*)"', body, re.I)
        title = html_module.unescape(img_m.group(1)) if img_m else f"Slide {idx + 1}"
        notes_m = re.search(r'<p class="notes">([^<]*)</p>', body, re.I)
        notes = html_module.unescape(notes_m.group(1)) if notes_m else ""
        scene = Scene(title=title, body=notes or title, narration=notes)
        scenes.append(scene)
        data_uri_m = re.search(r'src="(data:image/[^;]+;base64,[^"]+)"', body, re.I)
        if data_uri_m:
            png = _decode_data_uri_to_file(data_uri_m.group(1), work / f"slide-{idx:02d}.png")
            slide_images.append((png, scene))

    return LoadedDeck(
        topic=topic,
        scenes=scenes,
        slide_images=slide_images or None,
    )


def _extract_pptx_images(path: Path, work: Path) -> list[tuple[Path, int]]:
    Presentation, _ = _require_pptx()
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    images: list[tuple[Path, int]] = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    png_path = work / f"slide-{i:02d}.png"
                    png_path.write_bytes(shape.image.blob)
                    images.append((png_path, i))
                except (AttributeError, OSError):
                    pass
                break
    return images


def _load_pptx_deck(path: Path, work: Path) -> LoadedDeck:
    sidecar = _find_sidecar(path)
    images = _extract_pptx_images(path, work)

    if sidecar:
        deck = _load_json_deck(sidecar)
        if images:
            slide_images: list[tuple[Path, Scene]] = []
            for png_path, idx in images:
                if idx < len(deck.scenes):
                    deck.scenes[idx].slide_image = str(png_path)
                    slide_images.append((png_path, deck.scenes[idx]))
                else:
                    scene = Scene(title=f"Slide {idx + 1}", body=f"Slide {idx + 1}")
                    scene.slide_image = str(png_path)
                    deck.scenes.append(scene)
                    slide_images.append((png_path, scene))
            deck.slide_images = slide_images or None
        return deck

    Presentation, _ = _require_pptx()

    prs = Presentation(str(path))
    scenes: list[Scene] = []
    slide_images: list[tuple[Path, Scene]] = []

    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        png_path: Path | None = None
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunk = (shape.text or "").strip()
                if chunk:
                    texts.append(chunk)
        for png_path_c, idx in images:
            if idx == i:
                png_path = png_path_c
                break
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        title = texts[0] if texts else f"Slide {i + 1}"
        body = "\n".join(texts[1:]).strip() if len(texts) > 1 else ""
        scene = Scene(title=title, body=body or title, narration=notes)
        scenes.append(scene)
        if png_path and png_path.is_file():
            scene.slide_image = str(png_path)
            slide_images.append((png_path, scene))

    if not scenes:
        raise SystemExit(f"No slides found in {path.name}")
    topic = path.stem
    return LoadedDeck(topic=topic, scenes=scenes, slide_images=slide_images or None)


def _require_pymupdf():
    try:
        import fitz  # noqa: F401
    except ImportError as exc:
        raise SystemExit(
            "PyMuPDF is required to convert PDF slide decks.\n"
            "Install: pip install pymupdf\n"
            "Or: pip install 'arka-agent[drawings]'\n"
            "Tip: place a .meta.json sidecar beside the PDF for text-only conversion."
        ) from exc


def _load_pdf_deck(path: Path, work: Path) -> LoadedDeck:
    sidecar = _find_sidecar(path)
    if sidecar:
        return _load_json_deck(sidecar)

    _require_pymupdf()
    import fitz

    doc = fitz.open(path)
    try:
        if len(doc) == 0:
            raise SystemExit(f"PDF has no pages: {path.name}")
        scenes: list[Scene] = []
        slide_images: list[tuple[Path, Scene]] = []
        zoom = float(_env("SLIDES_PDF_ZOOM", "2.0"))
        matrix = fitz.Matrix(zoom, zoom)
        for i in range(len(doc)):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            png = work / f"slide-{i:02d}.png"
            png.write_bytes(pix.tobytes("png"))
            scene = Scene(title=f"Slide {i + 1}", body=f"Slide {i + 1}")
            scene.slide_image = str(png)
            scenes.append(scene)
            slide_images.append((png, scene))
        return LoadedDeck(topic=path.stem, scenes=scenes, slide_images=slide_images)
    finally:
        doc.close()


def _soffice_convert(input_path: Path, to_ext: str, out_dir: Path) -> Path:
    soffice = _which("soffice") or _which("libreoffice")
    if not soffice:
        raise SystemExit(
            f"LibreOffice (soffice) is required for direct {input_path.suffix} → .{to_ext} conversion.\n"
            "Install LibreOffice, or convert via JSON/HTML sidecar first."
        )
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        to_ext,
        "--outdir",
        str(out_dir),
        str(input_path),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        detail = (proc.stderr or proc.stdout or "").strip()
        raise SystemExit(f"LibreOffice conversion failed: {detail or proc.returncode}")
    out = out_dir / f"{input_path.stem}.{to_ext}"
    if not out.is_file():
        matches = sorted(out_dir.glob(f"{input_path.stem}*.{to_ext}"))
        if not matches:
            raise SystemExit(f"LibreOffice did not produce .{to_ext} output for {input_path.name}")
        out = matches[0]
    return out


def load_deck(path: Path, from_fmt: str | None = None, *, work: Path | None = None) -> LoadedDeck:
    path = path.expanduser().resolve()
    if not path.is_file():
        raise SystemExit(f"Input not found: {path}")
    fmt = detect_format(path, from_fmt)
    work_dir = work or Path(tempfile.mkdtemp(prefix="arka-slides-load-"))

    loaders = {
        "json": lambda p: _load_json_deck(p),
        "md": lambda p: _load_markdown_deck(p),
        "html": lambda p: _load_html_deck(p, work_dir),
        "pptx": lambda p: _load_pptx_deck(p, work_dir),
        "pdf": lambda p: _load_pdf_deck(p, work_dir),
    }
    return loaders[fmt](path)


def _render_convert_slides(
    scenes: list[Scene],
    *,
    work: Path,
    cfg: VideoConfig,
    existing_images: list[tuple[Path, Scene]] | None = None,
) -> list[tuple[Path, Scene]]:
    from arka.media.chart_slide import render_scene_visual, render_title_slide, scene_has_chart_visual

    import shutil

    existing = {i: img for i, (img, _) in enumerate(existing_images or [])}
    slide_images: list[tuple[Path, Scene]] = []
    for i, scene in enumerate(scenes):
        if i in existing and existing[i].is_file():
            slide_images.append((existing[i], scene))
            continue
        if scene.slide_image:
            src = Path(scene.slide_image).expanduser()
            if src.is_file():
                dest = work / f"slide-{i:02d}.png"
                shutil.copy2(src, dest)
                slide_images.append((dest, scene))
                continue
        out = work / f"slide-{i:02d}.png"
        if scene_has_chart_visual(scene):
            png = render_scene_visual(scene, work, cfg, index=i)
        else:
            png = render_title_slide(scene, out, cfg)
        slide_images.append((png, scene))
    return slide_images


def convert_deck(
    input_path: Path,
    *,
    output: Path,
    formats: list[str],
    from_fmt: str | None = None,
    cfg: VideoConfig | None = None,
) -> ExportBatch:
    """Convert an existing slide deck between supported formats."""
    input_path = input_path.expanduser()
    if not formats:
        raise SystemExit("No output formats specified.")

    src_fmt = detect_format(input_path, from_fmt)
    direct_pairs = {("pptx", "pdf"), ("pdf", "pptx")}
    if len(formats) == 1 and (src_fmt, formats[0]) in direct_pairs:
        try:
            out_dir = output.parent if output.suffix else output.parent
            out_dir.mkdir(parents=True, exist_ok=True)
            direct = _soffice_convert(input_path.resolve(), formats[0], out_dir)
            target = output if output.suffix else output.with_suffix(FORMAT_EXTENSIONS[formats[0]])
            if direct.resolve() != target.resolve():
                target.parent.mkdir(parents=True, exist_ok=True)
                direct.replace(target)
                direct = target
            fmt = formats[0]
            validator = _export_validator(fmt)
            if validator:
                try:
                    validator(direct)
                except Exception as exc:
                    _remove_invalid_export(direct)
                    raise SystemExit(f"{fmt} conversion produced an invalid file: {exc}") from exc
            return ExportBatch(saved=[direct], outputs={fmt: str(direct)})
        except SystemExit:
            if formats[0] not in SLIDE_FORMATS:
                raise
            print("  Direct office conversion unavailable — using scene pipeline …", file=sys.stderr)

    cfg = cfg or load_config()
    work = Path(tempfile.mkdtemp(prefix="arka-slides-convert-"))
    try:
        deck = load_deck(input_path, from_fmt, work=work)
        if not deck.scenes:
            raise SystemExit("No slides to convert.")

        print(
            f"Converting {input_path.name} ({src_fmt}) → {', '.join(formats)} ({len(deck.scenes)} slides)",
            file=sys.stderr,
        )

        slide_images = _render_convert_slides(
            deck.scenes,
            work=work,
            cfg=cfg,
            existing_images=deck.slide_images,
        )

        output_paths = _output_paths(output, formats)
        batch = _export_formats(
            slide_images,
            output_paths,
            topic=deck.topic,
            scenes=deck.scenes,
            cfg=cfg,
            credits=deck.credits,
        )

        if "json" not in batch.outputs:
            sidecar = output.with_suffix(".meta.json")
            sidecar.write_text(
                json.dumps(
                    _metadata_payload(
                        topic=deck.topic,
                        scenes=deck.scenes,
                        credits=deck.credits,
                        outputs=batch.outputs,
                    ),
                    indent=2,
                )
                + "\n",
                encoding="utf-8",
            )
        return batch
    finally:
        import shutil

        shutil.rmtree(work, ignore_errors=True)


def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Compose and convert presentation slide decks — pptx, pdf, html, md, json"
    )
    sub = p.add_subparsers(dest="cmd")

    p_compose = sub.add_parser("compose", help="Build slides from topic or script")
    p_compose.add_argument("--topic", help="Presentation topic (template or --llm script)")
    p_compose.add_argument("--script", help="JSON script file or inline JSON")
    p_compose.add_argument("--llm", action="store_true", help="Generate script via LLM")
    p_compose.add_argument(
        "--style",
        choices=SLIDE_STYLES,
        default=None,
        help="Deck tone: executive (default), academic, or pitch",
    )
    p_compose.add_argument(
        "--scenes",
        type=int,
        default=None,
        metavar="N",
        help="Scene count for --llm (default: let the LLM choose)",
    )
    p_compose.add_argument(
        "-f",
        "--format",
        default="auto",
        help="Export format: pptx, pdf, html, md, json, auto (default), or all",
    )
    p_compose.add_argument("-o", "--output", help="Output path (extension adjusted per --format)")
    p_compose.set_defaults(func=cmd_compose)

    p_convert = sub.add_parser("convert", help="Convert an existing deck between formats")
    p_convert.add_argument("input", help="Input deck (.pptx, .pdf, .html, .md, .json)")
    p_convert.add_argument(
        "-f",
        "--format",
        "--to",
        dest="format",
        default="auto",
        help="Output format: pptx, pdf, html, md, json, auto (from -o extension), or all",
    )
    p_convert.add_argument(
        "--from",
        dest="from_fmt",
        default=None,
        help="Input format (auto-detected from extension when omitted)",
    )
    p_convert.add_argument("-o", "--output", help="Output path (extension adjusted per --format)")
    p_convert.set_defaults(func=cmd_convert)

    p_parse = sub.add_parser("parse", help="Parse natural language → compose args")
    p_parse.add_argument("text", nargs="+")
    p_parse.set_defaults(func=cmd_parse)

    p_check = sub.add_parser("check", help="Verify exporters, Pillow, stock photo keys")
    p_check.set_defaults(func=cmd_check)

    return p


def cmd_check(_args: argparse.Namespace) -> int:
    ok = True
    try:
        _require_pptx()
        print("✓ pptx (python-pptx)")
    except SystemExit:
        print("✗ pptx — pip install python-pptx", file=sys.stderr)
        ok = False
    try:
        from arka.media.compose_video import _require_pillow

        _require_pillow()
        print("✓ pdf (Pillow)")
    except SystemExit:
        print("✗ pdf — pip install Pillow", file=sys.stderr)
        ok = False
    print("✓ html (stdlib)")
    print("✓ md (stdlib)")
    print("✓ json (stdlib)")
    try:
        _require_pymupdf()
        print("✓ pdf input (pymupdf)")
    except SystemExit:
        print("  pdf input — optional: pip install pymupdf (or use .meta.json sidecar)")
    if _which("soffice") or _which("libreoffice"):
        print("✓ LibreOffice (direct pptx↔pdf)")
    else:
        print("  LibreOffice optional for direct pptx↔pdf")
    if any_source_available():
        print("✓ Stock photo source configured")
    else:
        print(f"  {stock_setup_hint('compose_slides')}")
    cfg = load_config()
    print(f"  Slide size: {cfg.width}x{cfg.height}")
    return 0 if ok else 1


def cmd_parse(args: argparse.Namespace) -> int:
    argv = nl_to_argv(" ".join(args.text))
    if not argv:
        return 1
    print(" ".join(shlex.quote(a) for a in argv))
    return 0


def _report_export_batch(batch: ExportBatch, *, saved_label: str = "Saved slides") -> Path | None:
    from arka.core.output import user_msg

    open_path: Path | None = None
    if batch.failed.get("pptx"):
        user_msg("✗ Could not create a valid PowerPoint file.")
        if batch.fallback_md:
            user_msg(f"  Saved outline instead: {batch.fallback_md}")
            open_path = batch.fallback_md
    for path in batch.saved:
        if batch.fallback_md and path == batch.fallback_md and batch.failed.get("pptx"):
            continue
        print(f"{saved_label}: {path}")
        if open_path is None:
            open_path = path
    return open_path


def cmd_compose(args: argparse.Namespace) -> int:
    from arka.core.output import debug_msg, user_msg

    scenes: list[Scene] = []
    topic = extract_slides_topic((args.topic or "").strip()) or normalize_topic((args.topic or "").strip())
    style = normalize_slide_style(args.style)

    if args.script:
        raw = Path(args.script).expanduser().read_text(encoding="utf-8") if Path(args.script).is_file() else args.script
        scenes = _enrich_slide_scenes(_parse_scenes_json(raw))
        if not topic and scenes:
            topic = scenes[0].title

    if not scenes and topic:
        mode = _slides_script_mode(args)
        if mode == "llm":
            if args.scenes is not None:
                user_msg(f"Writing {style} slide script ({args.scenes} slides) …")
            else:
                user_msg(f"Writing {style} slide script …")
            try:
                scenes = _llm_slides_script(topic, scenes=args.scenes, style=style)
                if _slides_script_needs_shortening(scenes):
                    user_msg("Deck too dense — tightening …")
                    scenes = _llm_slides_summarize_script(topic, scenes, style=style)
                if args.scenes is None:
                    debug_msg(f"LLM chose {len(scenes)} slides")
            except SystemExit as exc:
                user_msg(f"LLM script failed ({exc}); using {style} template.")
                scenes = _template_slides_script(topic, style=style)
        else:
            scenes = _template_slides_script(topic, style=style)

    if not scenes:
        user_msg("Provide --topic or --script")
        return 1
    if not topic:
        topic = scenes[0].title

    if _llm_available() and any(not s.image_keywords for s in scenes):
        debug_msg("Choosing stock photo keywords with LLM …")
        scenes = _llm_slides_enrich_image_keywords(topic, scenes)

    label = topic_label(topic)
    debug_msg(f"Topic: {label} ({style})")
    try:
        formats = parse_formats_arg(args.format)
    except SystemExit as exc:
        user_msg(str(exc))
        return 1
    primary_fmt = formats[0]
    out = (
        Path(args.output).expanduser()
        if args.output
        else _default_output(topic, primary_fmt if len(formats) == 1 else "pptx")
    )
    fmt_label = ", ".join(formats)
    user_msg(f"Composing {len(scenes)} slides ({fmt_label}) …")
    batch = compose(scenes, output=out, topic=topic, formats=formats)
    open_path = _report_export_batch(batch)
    meta = out.with_suffix(".meta.json")
    if meta.is_file():
        debug_msg(f"Metadata: {meta}")
    elif "json" in {p.suffix.lstrip(".") for p in batch.saved}:
        json_out = next(p for p in batch.saved if p.suffix == ".json")
        debug_msg(f"Metadata: {json_out}")
    if open_path and _env("OPEN_SLIDES", "1") not in {"0", "false"}:
        _open_slides(open_path)
    return 0 if batch.saved else 1


def _format_from_output_path(output: Path) -> str | None:
    ext = output.suffix.lstrip(".").lower()
    return normalize_format(ext) or None


def cmd_convert(args: argparse.Namespace) -> int:
    input_path = Path(args.input).expanduser()
    output_arg = Path(args.output).expanduser() if args.output else None

    fmt_raw = args.format
    if fmt_raw in {None, "", "auto"} and output_arg and output_arg.suffix:
        detected = _format_from_output_path(output_arg)
        if detected:
            fmt_raw = detected

    try:
        formats = parse_formats_arg(fmt_raw)
    except SystemExit as exc:
        print(exc, file=sys.stderr)
        return 1

    if output_arg:
        output = output_arg
    else:
        stem = input_path.with_suffix("")
        output = (
            stem.with_suffix(FORMAT_EXTENSIONS[formats[0]])
            if len(formats) == 1
            else stem
        )

    try:
        batch = convert_deck(
            input_path,
            output=output,
            formats=formats,
            from_fmt=args.from_fmt,
        )
    except SystemExit as exc:
        print(exc, file=sys.stderr)
        return 1

    open_path = _report_export_batch(batch, saved_label="Saved")
    meta = output.with_suffix(".meta.json")
    if meta.is_file():
        print(f"Metadata: {meta}")
    elif "json" in {p.suffix.lstrip(".") for p in batch.saved}:
        json_out = next(p for p in batch.saved if p.suffix == ".json")
        print(f"Metadata: {json_out}")
    if open_path and _env("OPEN_SLIDES", "1") not in {"0", "false"}:
        _open_slides(open_path)
    return 0 if batch.saved else 1


def _open_slides(path: Path) -> None:
    if sys.platform == "darwin":
        subprocess.Popen(["open", str(path)], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    elif _which("xdg-open"):
        subprocess.Popen(["xdg-open", str(path)], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)


_COMPOSE_SUBCMDS = frozenset({"compose", "convert", "parse", "check"})
_COMPOSE_FLAG_PREFIXES = ("--topic", "--script", "--llm", "--style", "--scenes", "--format", "--output")


def _normalize_compose_argv(argv: list[str]) -> list[str]:
    """Accept legacy ``--topic`` argv (no ``compose`` subcommand) and NL phrases."""
    # Some command dispatchers use ``--`` to separate their own arguments from
    # the command they launch.  It is not part of compose_slides' CLI, so avoid
    # treating the following valid subcommand as natural-language input.
    while len(argv) > 1 and argv[0] == "--":
        argv = argv[1:]
    if not argv:
        return argv
    head = argv[0]
    if head in _COMPOSE_SUBCMDS or head in {"-h", "--help"}:
        return argv
    if head.startswith("-") and any(
        a == flag or a.startswith(f"{flag}=") for a in argv for flag in _COMPOSE_FLAG_PREFIXES
    ):
        return ["compose", *argv]
    nl = nl_to_argv(" ".join(argv))
    return nl if nl else argv


def main(argv: list[str] | None = None) -> int:
    from arka.env import load_env

    load_env()
    argv = list(argv if argv is not None else sys.argv[1:])
    if not argv:
        build_parser().print_help()
        return 0
    argv = _normalize_compose_argv(argv)
    if argv[0] not in _COMPOSE_SUBCMDS and argv[0] not in {"-h", "--help"}:
        build_parser().print_help()
        return 1
    parser = build_parser()
    args = parser.parse_args(argv)
    if not getattr(args, "func", None):
        parser.print_help()
        return 0
    return int(args.func(args))


if __name__ == "__main__":
    raise SystemExit(main())
