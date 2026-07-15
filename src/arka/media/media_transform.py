"""Content-aware transformations between media types."""
from __future__ import annotations

import argparse
import json
import re
from pathlib import Path


def _source_text(source: str) -> tuple[str, str]:
    path = Path(source).expanduser()
    if path.is_file():
        from arka.youtube.transcript import try_transcript_for_media
        result = try_transcript_for_media(path)
        if result:
            return result[1], f"media:{path.name}"
        if path.suffix.lower() in {".txt", ".md", ".html"}:
            return re.sub(r"<[^>]+>", " ", path.read_text(encoding="utf-8", errors="replace")), f"file:{path.name}"
        if path.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}:
            try:
                import pytesseract
                from PIL import Image
                return pytesseract.image_to_string(Image.open(path)), f"image:{path.name}"
            except ImportError as exc:
                raise RuntimeError("Image-to-book requires OCR dependencies: pip install pytesseract Pillow") from exc
        if path.suffix.lower() == ".pptx":
            try:
                from pptx import Presentation
                text = []
                for index, slide in enumerate(Presentation(str(path)).slides, 1):
                    parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    text.append(f"## Slide {index}\n\n" + "\n".join(parts))
                return "\n\n".join(text), f"slides:{path.name}"
            except ImportError as exc:
                raise RuntimeError("PPTX-to-book requires python-pptx") from exc
        if path.suffix.lower() == ".pdf":
            try:
                import fitz
                return "\n\n".join(page.get_text() for page in fitz.open(path)), f"pdf:{path.name}"
            except ImportError as exc:
                raise RuntimeError("PDF-to-book requires PyMuPDF") from exc
        raise RuntimeError("Could not extract text from this media; install ffmpeg/Whisper or provide captions")
    if source.startswith(("http://", "https://")):
        if "list=" in source or "playlist" in source.lower():
            from arka.media.batch import _playlist_entries
            from arka.youtube.transcript import get_transcript
            parts: list[str] = []
            for video_id, title in _playlist_entries(source, 50):
                try:
                    _id, text = get_transcript(video_id, label=title)
                    parts.append(f"## {title}\n\n{text}")
                except Exception:
                    continue
            if parts:
                return "\n\n".join(parts), source
        from arka.youtube.transcript import get_transcript
        _video_id, text = get_transcript(source)
        return text, source
    raise RuntimeError("source must be a local media/text file or a YouTube URL")


def to_book(source: str, output: Path) -> Path:
    text, label = _source_text(source)
    paragraphs = [part.strip() for part in re.split(r"\n\s*\n|(?<=[.!?])\s+(?=[A-Z])", text) if part.strip()]
    chapters = [paragraphs[index : index + 12] for index in range(0, len(paragraphs), 12)] or [[]]
    lines = [f"# {Path(source).stem.replace('_', ' ').title() or 'Arka Book'}", "", f"_Source: {label}_", ""]
    for index, chapter in enumerate(chapters, 1):
        lines += [f"## Chapter {index}", "", "\n\n".join(chapter), ""]
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(lines), encoding="utf-8")
    return output


def to_podcast(source: str, output: Path) -> Path:
    text, _label = _source_text(source)
    try:
        from arka.voice.edge_speak import synthesize_to_file
    except ImportError as exc:
        raise RuntimeError("Install edge-tts for podcast output: arka tts-setup") from exc
    return synthesize_to_file(text, output)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(prog="arka media transform")
    parser.add_argument("source")
    parser.add_argument("--to", choices=("book", "podcast"), required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args(argv)
    output = Path(args.output).expanduser()
    result = to_book(args.source, output) if args.to == "book" else to_podcast(args.source, output)
    print(json.dumps({"source": args.source, "target": args.to, "output": str(result)}))
    return 0
